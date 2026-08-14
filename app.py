import os
import io
import re
import json
import time
import math
import base64
import hashlib
import shutil
import threading
import zipfile
import gc
import tempfile
from copy import deepcopy, copy
from dataclasses import dataclass
from typing import List, Tuple, Optional
from datetime import datetime, timedelta
from difflib import SequenceMatcher

import requests
import streamlit as st
from PIL import Image, ImageOps, ImageEnhance, ImageFilter, ImageStat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor

try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
except Exception:
    pass


try:
    import pytesseract
except Exception:
    pytesseract = None

try:
    from paddleocr import PaddleOCR
except Exception:
    PaddleOCR = None

try:
    import openpyxl
except Exception:
    openpyxl = None

_PADDLE_OCR = None


BASE_DIR = os.path.dirname(os.path.abspath(__file__))

TEMPLATE_PPT = os.path.join(BASE_DIR, "templates", "sample_template.pptx")
DAILY_TEMPLATE_PPT = os.path.join(BASE_DIR, "templates", "sample_template2.pptx")
BAD_PHOTO_DIR = os.path.join(BASE_DIR, "bad_photo_upload")
SHARED_NOTICE_FILE = os.path.join(BASE_DIR, "shared_notice.md")
SHARED_NOTICE_META_FILE = os.path.join(BASE_DIR, "shared_notice_meta.json")

BASE_FONT_SIZE_PT = 35
OUTPUT_PPT_NAME = "TBM_완성본.pptx"
DAILY_OUTPUT_PPT_NAME = "일일안전회의_완성본.pptx"
APP_VERSION = "26년 8월 일일안전회의 고속분류 버전"

# 대량 업로드/고용량 사진 안정화 설정
TBM_IMAGE_MAX_SIZE = 1200
TBM_IMAGE_QUALITY = 82
DAILY_IMAGE_MAX_SIZE = 1400
DAILY_IMAGE_QUALITY = 84
MAX_TBM_FILES_SOFT_WARN = 35
MAX_DAILY_FILES_SOFT_WARN = 35
TRANSLATION_BATCH_SIZE = 15

PHOTO_BOX_TEXT = "PHOTO_BOX"
KO_BOX_TEXT = "1"
ZH_BOX_TEXT = "2"
VI_BOX_TEXT = "3"
MY_BOX_TEXT = "4"

DATE_BOX_TEXT = "DATE_BOX"

DAILY_PHOTO_BOX_TEXT = "PHOTO_BOX_1"
DAILY_TEXT_BOX_TEXT = "TEXT_BOX_1"
TIME_BOX_TEXT = "TIME_BOX_1"
HOLD_POINT_TEXT = "HOLD POINT"

BAD_PHOTO_LIMIT_MB = 300
BAD_PHOTO_LIMIT_BYTES = BAD_PHOTO_LIMIT_MB * 1024 * 1024
BAD_PHOTO_EXPIRE_SECONDS = 24 * 60 * 60



@dataclass
class SlideData:
    image_path: str
    ko: str = ""
    zh: str = ""
    vi: str = ""
    my: str = ""
    orientation: str = ""


@dataclass
class DailySlideData:
    image_path: str
    text: str = ""
    orientation: str = ""


@dataclass
class MaterialWorkItem:
    image_path: str
    original_name: str = ""
    upload_index: int = 0
    ocr_text: str = ""
    work_type: str = "material"
    company: str = "기타업체"
    number: int = 0


COMPANY_ORDER = [
    "원영건업",
    "청암기업",
    "유셀네트웍스",
    "엠케이지",
    "금성",
    "웰시스템",
    "KEC",
    "청오",
    "우신에이스",
    "MS건설",
    "진솔",
    "장한건설",
    "신영기초개발",
    "KCC",
    "씨즌텍",
]

COMPANY_ALIAS = {
    "원영건업": ["원영건업", "원영"],
    "청암기업": ["청암기업", "청암"],
    "유셀네트웍스": ["유셀네트웍스", "유셀네트윅스", "유셀네트", "유셀"],
    "엠케이지": ["엠케이지", "MKG", "mkg"],
    "KEC": ["KEC", "kec", "케이이씨", "케이씨", "케이"],
    "우신에이스": ["우신에이스", "우신"],
    "진솔": ["진솔"],
    "장한건설": ["장한건설", "장한"],
}

HIGH_RISK_KEYWORDS = [
    "25대 고위험",
    "25대고위험",
    "25 대 고위험",
    "25대",
    "고위험",
    "고 위험",
]


def hide_streamlit_ui():
    st.markdown(
        """
        <style>
        #MainMenu {visibility: hidden;}
        header {visibility: hidden;}
        footer {visibility: hidden;}
        .stDeployButton {display: none;}
        [data-testid="stToolbar"] {display: none;}
        [data-testid="stDecoration"] {display: none;}
        [data-testid="stStatusWidget"] {display: none;}
        [data-testid="manage-app-button"] {display: none;}

        .block-container {
            padding-top: 1.2rem !important;
        }

        hr {
            margin-top: 0.35rem !important;
            margin-bottom: 0.35rem !important;
        }

        h2, h3, h4 {
            margin-top: 0.35rem !important;
            margin-bottom: 0.35rem !important;
        }

        div[data-testid="stMarkdownContainer"] p {
            margin-bottom: 0.25rem !important;
        }

        div[data-testid="stFileUploader"] {
            margin-top: -0.2rem !important;
            margin-bottom: 0.35rem !important;
        }

        div[data-testid="stVerticalBlock"] {
            gap: 0.35rem !important;
        }

        .bad-photo-small-button button {
            min-height: 34px !important;
            padding: 0.25rem 0.4rem !important;
            font-size: 12px !important;
        }

        /* 일일안전회의 생성 버튼이 모바일/카카오 인앱브라우저에서
           다른 영역에 덮이지 않도록 독립 레이어로 유지 */
        div[data-testid="stButton"] {
            position: relative !important;
            z-index: 20 !important;
            clear: both !important;
        }

        div[data-testid="stDownloadButton"] {
            position: relative !important;
            z-index: 20 !important;
            clear: both !important;
        }

        .block-container {
            padding-bottom: 7rem !important;
        }

        @media (max-width: 768px) {
            div[data-testid="stButton"] button {
                min-height: 52px !important;
                font-size: 17px !important;
                font-weight: 700 !important;
            }

            div[data-testid="stDownloadButton"] button {
                min-height: 48px !important;
            }

            div[data-testid="stFileUploader"] {
                margin-bottom: 0.8rem !important;
            }

            .block-container {
                padding-left: 0.9rem !important;
                padding-right: 0.9rem !important;
                padding-bottom: 9rem !important;
            }
        }
        </style>
        """,
        unsafe_allow_html=True
    )


def render_app_title():
    st.markdown(
        f"""
        <div style="margin-top:-30px; margin-bottom:8px;">
            <h2 style="
                font-size:24px;
                font-weight:700;
                margin:0;
                padding:0;
                line-height:1.2;
            ">
                🚧 TBM 교육자료 자동 번역 생성기 [{APP_VERSION}]
            </h2>
        </div>
        """,
        unsafe_allow_html=True
    )


def safe_filename(filename: str) -> str:
    name = os.path.basename(filename)
    name = re.sub(r"[^a-zA-Z0-9가-힣._ -]", "_", name)
    return name[:120]


def format_size(size_bytes: int) -> str:
    if size_bytes < 1024:
        return f"{size_bytes}B"
    if size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f}KB"
    return f"{size_bytes / 1024 / 1024:.1f}MB"


def ensure_bad_photo_dir():
    os.makedirs(BAD_PHOTO_DIR, exist_ok=True)


def cleanup_old_bad_photo_files():
    ensure_bad_photo_dir()
    now = time.time()

    for name in os.listdir(BAD_PHOTO_DIR):
        path = os.path.join(BAD_PHOTO_DIR, name)
        if os.path.isfile(path):
            if now - os.path.getmtime(path) > BAD_PHOTO_EXPIRE_SECONDS:
                try:
                    os.remove(path)
                except Exception:
                    pass


def get_bad_photo_size() -> int:
    ensure_bad_photo_dir()
    total = 0

    for name in os.listdir(BAD_PHOTO_DIR):
        path = os.path.join(BAD_PHOTO_DIR, name)
        if os.path.isfile(path):
            total += os.path.getsize(path)

    return total


def save_bad_photo_file(uploaded_file):
    ensure_bad_photo_dir()

    file_bytes = uploaded_file.getvalue()
    file_size = len(file_bytes)
    original_filename = safe_filename(uploaded_file.name)

    for existing_name in os.listdir(BAD_PHOTO_DIR):
        if existing_name.endswith(original_filename):
            existing_path = os.path.join(BAD_PHOTO_DIR, existing_name)
            if os.path.isfile(existing_path) and os.path.getsize(existing_path) == file_size:
                return existing_path

    current_size = get_bad_photo_size()

    if current_size + file_size > BAD_PHOTO_LIMIT_BYTES:
        raise ValueError(
            f"부적합사진 용량 초과: 현재 {format_size(current_size)} / "
            f"추가 {format_size(file_size)} / 최대 {BAD_PHOTO_LIMIT_MB}MB"
        )

    timestamp = time.strftime("%Y%m%d_%H%M%S")
    save_name = f"{timestamp}_{original_filename}"
    save_path = os.path.join(BAD_PHOTO_DIR, save_name)

    with open(save_path, "wb") as f:
        f.write(file_bytes)

    return save_path


def save_generated_ppt_to_bad_photo_storage(ppt_bytes: io.BytesIO, filename: str):
    ensure_bad_photo_dir()

    ppt_bytes.seek(0)
    data = ppt_bytes.getvalue()
    file_size = len(data)

    current_size = get_bad_photo_size()

    if current_size + file_size > BAD_PHOTO_LIMIT_BYTES:
        ppt_bytes.seek(0)
        return False

    safe_name = safe_filename(filename)
    timestamp = time.strftime("%Y%m%d_%H%M%S")
    save_name = f"{timestamp}_{safe_name}"
    save_path = os.path.join(BAD_PHOTO_DIR, save_name)

    with open(save_path, "wb") as f:
        f.write(data)

    ppt_bytes.seek(0)
    return True


def make_bad_photo_zip():
    ensure_bad_photo_dir()

    zip_buffer = io.BytesIO()

    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
        for name in sorted(os.listdir(BAD_PHOTO_DIR)):
            path = os.path.join(BAD_PHOTO_DIR, name)
            if os.path.isfile(path):
                zip_file.write(path, arcname=name)

    zip_buffer.seek(0)
    return zip_buffer


def delete_all_bad_photo_files():
    ensure_bad_photo_dir()
    deleted_count = 0

    for name in os.listdir(BAD_PHOTO_DIR):
        path = os.path.join(BAD_PHOTO_DIR, name)
        if os.path.isfile(path):
            try:
                os.remove(path)
                deleted_count += 1
            except Exception:
                pass

    return deleted_count


def make_thumbnail_bytes(path: str, max_size: int = 180):
    try:
        img = Image.open(path)

        try:
            img.seek(0)
        except Exception:
            pass

        try:
            img = ImageOps.exif_transpose(img)
        except Exception:
            pass

        if img.mode != "RGB":
            img = img.convert("RGB")

        img.thumbnail((max_size, max_size))

        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        buf.seek(0)
        return buf

    except Exception:
        return None


def delete_file(path: str):
    try:
        if os.path.exists(path) and os.path.isfile(path):
            os.remove(path)
            return True
    except Exception:
        pass
    return False


def render_bad_photo_storage():
    cleanup_old_bad_photo_files()

    files = []
    ensure_bad_photo_dir()

    for name in sorted(os.listdir(BAD_PHOTO_DIR), reverse=True):
        path = os.path.join(BAD_PHOTO_DIR, name)
        if os.path.isfile(path):
            files.append((name, path, os.path.getsize(path)))

    col_title, col_zip, col_delete = st.columns([5.2, 0.9, 0.9])

    with col_title:
        used = get_bad_photo_size()
        st.markdown(
            f"""
            <div style="
                display:flex;
                align-items:center;
                gap:12px;
                margin:0;
                padding:0;
                line-height:1.1;
                min-height:34px;
            ">
                <span style="font-size:48px; font-weight:800; color:#006400;">부적합사진</span>
                <span style="font-size:13px; color:#666;">
                    용량 {format_size(used)} / {BAD_PHOTO_LIMIT_MB}MB
                </span>
            </div>
            """,
            unsafe_allow_html=True
        )

    with col_zip:
        st.markdown("<div class='bad-photo-small-button'>", unsafe_allow_html=True)
        if files:
            st.download_button(
                "ZIP",
                data=make_bad_photo_zip(),
                file_name="부적합사진_전체.zip",
                mime="application/zip",
                use_container_width=True,
                key="bad_photo_download_all_zip"
            )
        else:
            st.button("ZIP", disabled=True, use_container_width=True, key="bad_photo_download_all_zip_disabled")
        st.markdown("</div>", unsafe_allow_html=True)

    with col_delete:
        st.markdown("<div class='bad-photo-small-button'>", unsafe_allow_html=True)
        if st.button("전체삭제", use_container_width=True, key="bad_photo_delete_all"):
            delete_all_bad_photo_files()
            st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)

    upload_files = st.file_uploader(
        "부적합사진 파일 업로드",
        accept_multiple_files=True,
        type=None,  # 확장자 제한 없음: 사진, 영상, MP3, 압축파일, 문서 등 업로드 가능
        key="daily_bad_photo_upload_uploader",
        help="확장자 제한 없음: 사진, 영상, MP3, 압축파일, 문서 등 거의 모든 파일 업로드 가능"
    )

    if upload_files:
        uploaded_count = 0

        for file in upload_files:
            try:
                save_bad_photo_file(file)
                uploaded_count += 1
            except Exception as e:
                st.error(str(e))

        if uploaded_count > 0:
            st.success(f"{uploaded_count}개 등록 완료")
            st.rerun()

    if files:
        with st.expander(f"파일 목록 보기 ({len(files)}개)", expanded=False):
            for idx, (name, path, size) in enumerate(files, start=1):
                col1, col2, col3 = st.columns([1, 3, 0.5])

                thumb = make_thumbnail_bytes(path)

                with col1:
                    if thumb:
                        st.image(thumb, width=110)
                    else:
                        st.write("파일")
                        st.caption("미리보기 불가")

                with col2:
                    st.caption(name)
                    st.caption(format_size(size))

                    with open(path, "rb") as f:
                        st.download_button(
                            label="다운로드",
                            data=f,
                            file_name=name,
                            use_container_width=True,
                            key=f"bad_photo_download_{name}_{idx}"
                        )

                with col3:
                    if st.button("X", key=f"bad_photo_delete_{name}_{idx}", use_container_width=True):
                        if delete_file(path):
                            st.rerun()
                        else:
                            st.error("삭제 실패")
    else:
        st.info("부적합사진 파일 없음.")


def get_image_orientation(path: str) -> str:
    """EXIF 방향 보정 후 실제 이미지 방향을 반환."""
    try:
        img = Image.open(path)
        try:
            img.seek(0)
        except Exception:
            pass
        try:
            img = ImageOps.exif_transpose(img)
        except Exception:
            pass
        w, h = img.size
        if w > h:
            return "landscape"
        if h > w:
            return "portrait"
        return "square"
    except Exception:
        return "unknown"


def convert_to_jpg(
    input_path: str,
    max_size: int = TBM_IMAGE_MAX_SIZE,
    quality: int = TBM_IMAGE_QUALITY
) -> str:
    """
    PPT 삽입용 JPG 변환.
    - EXIF Orientation을 실제 픽셀 방향으로 반영한다.
    - 가로/세로 방향을 임의 변경하지 않는다.
    - 긴 변 기준으로 축소해 대량 사진/고용량 사진 생성 실패를 줄인다.
    """
    try:
        img = Image.open(input_path)

        try:
            img.seek(0)
        except Exception:
            pass

        original_orientation = get_image_orientation(input_path)

        try:
            img = ImageOps.exif_transpose(img)
        except Exception:
            pass

        if img.mode not in ("RGB", "L"):
            # 투명 PNG/WEBP는 흰 배경으로 합성해서 PPT 호환성 확보
            if "A" in img.getbands():
                bg = Image.new("RGB", img.size, (255, 255, 255))
                bg.paste(img, mask=img.getchannel("A"))
                img = bg
            else:
                img = img.convert("RGB")
        elif img.mode == "L":
            img = img.convert("RGB")

        width, height = img.size
        longest = max(width, height)

        if longest > max_size:
            ratio = max_size / longest
            new_width = max(1, int(width * ratio))
            new_height = max(1, int(height * ratio))
            img = img.resize((new_width, new_height), Image.LANCZOS)

        converted_orientation = "landscape" if img.width > img.height else "portrait" if img.height > img.width else "square"
        if original_orientation in ("landscape", "portrait") and converted_orientation != original_orientation:
            raise ValueError("이미지 변환 중 가로/세로 방향이 변경되었습니다.")

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        output_path = tmp.name
        tmp.close()

        img.save(output_path, format="JPEG", quality=quality, optimize=True, progressive=True)
        img.close()
        return output_path

    except Exception as e:
        raise ValueError(f"이미지 변환 실패: {e}")


def get_korean_date_text():
    weekdays = ["월요일", "화요일", "수요일", "목요일", "금요일", "토요일", "일요일"]
    now = datetime.now()
    return f"{now.year}년 {now.month:02d}월 {now.day:02d}일 {weekdays[now.weekday()]}"


def translate_batch_with_gpt(api_key: str, korean_list: List[str]):
    url = "https://api.openai.com/v1/responses"

    joined_text = "\n".join([f"{i+1}. {txt}" for i, txt in enumerate(korean_list)])

    prompt = f"""
다음 한국어 안전 문구들을 건설현장 TBM용으로 짧고 명확하게 번역하라.

조건:
- 반드시 JSON 배열만 출력
- 설명 금지
- 코드블록 금지
- 각 항목은 zh, vi, my 포함
- 입력 개수와 출력 개수는 반드시 같아야 함

입력:
{joined_text}

출력:
[
  {{
    "zh":"중국어",
    "vi":"베트남어",
    "my":"미얀마어"
  }}
]
"""

    headers = {
        "Authorization": f"Bearer {api_key.strip()}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": "gpt-4o-mini",
        "input": prompt
    }

    resp = requests.post(url, headers=headers, json=payload, timeout=60)

    if resp.status_code != 200:
        raise Exception(f"API Error: {resp.text}")

    data = resp.json()

    text = ""
    if "output_text" in data and data["output_text"]:
        text = data["output_text"]
    else:
        for item in data.get("output", []):
            for c in item.get("content", []):
                if c.get("type") == "output_text":
                    text += c.get("text", "")

    text = text.replace("```json", "").replace("```", "").strip()

    try:
        parsed = json.loads(text)
    except Exception:
        text = re.sub(r'[\x00-\x1F]+', ' ', text)
        parsed = json.loads(text)

    if not isinstance(parsed, list):
        raise ValueError("GPT 응답이 배열이 아닙니다.")

    if len(parsed) != len(korean_list):
        raise ValueError(
            f"번역 개수 불일치: 입력 {len(korean_list)} / 출력 {len(parsed)}"
        )

    for item in parsed:
        if not all(k in item for k in ("zh", "vi", "my")):
            raise ValueError("번역 결과에 zh, vi, my 키가 없습니다.")

    return parsed


def translate_all_with_gpt(api_key: str, korean_list: List[str]):
    """사진 수가 많을 때 API 응답 길이/시간초과를 줄이기 위해 나눠 번역."""
    results = []
    for start_idx in range(0, len(korean_list), TRANSLATION_BATCH_SIZE):
        chunk = korean_list[start_idx:start_idx + TRANSLATION_BATCH_SIZE]
        results.extend(translate_batch_with_gpt(api_key, chunk))
    return results


def iter_all_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in iter_all_shapes(shape.shapes):
                yield sub_shape


def has_text(shape):
    return hasattr(shape, "has_text_frame") and shape.has_text_frame


def normalize_text(text: str) -> str:
    return str(text).strip().replace("\n", "").replace("\r", "")


def normalize_for_match(text: str) -> str:
    return re.sub(r"\s+", "", str(text or "")).upper()


def get_paddle_ocr():
    """PaddleOCR 지연 로딩. 여러 버전 호환을 위해 옵션을 단계적으로 시도."""
    global _PADDLE_OCR

    if PaddleOCR is None:
        return None

    if _PADDLE_OCR is not None:
        return _PADDLE_OCR

    option_list = [
        dict(lang="korean", use_angle_cls=True, show_log=False, det_db_box_thresh=0.25, det_db_unclip_ratio=2.0),
        dict(lang="korean", use_angle_cls=True, show_log=False),
        dict(lang="korean", use_angle_cls=True),
        dict(lang="korean"),
    ]

    for opts in option_list:
        try:
            _PADDLE_OCR = PaddleOCR(**opts)
            return _PADDLE_OCR
        except Exception:
            continue

    _PADDLE_OCR = None
    return None


def normalize_ocr_text(text: str) -> str:
    text = str(text or "")
    text = text.replace("|", "I").replace("﹣", "-").replace("–", "-").replace("—", "-")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def resize_for_ocr(img: Image.Image, min_longest: int = 2600, max_longest: int = 4600) -> Image.Image:
    """OCR용 크기 보정. 작으면 키우고, 너무 크면 메모리 보호 수준에서만 줄임."""
    width, height = img.size
    longest = max(width, height)

    if longest < min_longest:
        ratio = min_longest / max(1, longest)
        img = img.resize((int(width * ratio), int(height * ratio)), Image.LANCZOS)
    elif longest > max_longest:
        ratio = max_longest / longest
        img = img.resize((int(width * ratio), int(height * ratio)), Image.LANCZOS)

    return img


def make_ocr_preprocess_variants(img: Image.Image) -> List[Image.Image]:
    """
    현장 문서용 강화 OCR 전처리.
    같은 영역에서 여러 이미지 버전을 만들어 OCR 결과를 합친다.
    느려지지만 작은 글씨/그림자/흐림 대응력이 올라간다.
    """
    try:
        img = ImageOps.exif_transpose(img)
    except Exception:
        pass

    if img.mode != "RGB":
        img = img.convert("RGB")

    img = resize_for_ocr(img)

    variants = []

    # 1) 원본 RGB 확대본
    variants.append(img)

    gray = img.convert("L")
    gray = ImageOps.autocontrast(gray)

    # 2) 기본 흑백 대비 강화
    v1 = ImageEnhance.Contrast(gray).enhance(1.8)
    v1 = ImageEnhance.Sharpness(v1).enhance(2.2)
    v1 = v1.filter(ImageFilter.SHARPEN)
    variants.append(v1)

    # 3) 더 강한 대비/샤프닝
    v2 = ImageEnhance.Contrast(gray).enhance(2.6)
    v2 = ImageEnhance.Sharpness(v2).enhance(3.0)
    v2 = v2.filter(ImageFilter.SHARPEN)
    variants.append(v2)

    # 4) 밝기 보정 + 대비
    v3 = ImageEnhance.Brightness(gray).enhance(1.12)
    v3 = ImageOps.autocontrast(v3)
    v3 = ImageEnhance.Contrast(v3).enhance(2.1)
    variants.append(v3)

    # 5) 이진화 1 - 일반 문서
    try:
        v4 = gray.point(lambda p: 255 if p > 165 else 0)
        variants.append(v4)
    except Exception:
        pass

    # 6) 이진화 2 - 어두운 사진/그림자 대응
    try:
        v5 = gray.point(lambda p: 255 if p > 135 else 0)
        variants.append(v5)
    except Exception:
        pass

    # 7) 가장 강한 확대 + 샤프닝
    try:
        w, h = img.size
        if max(w, h) < 4200:
            v6 = img.resize((int(w * 1.35), int(h * 1.35)), Image.LANCZOS).convert("L")
            v6 = ImageOps.autocontrast(v6)
            v6 = ImageEnhance.Contrast(v6).enhance(2.4)
            v6 = ImageEnhance.Sharpness(v6).enhance(3.5)
            variants.append(v6)
    except Exception:
        pass

    return variants


def ocr_with_paddle(img: Image.Image) -> str:
    ocr = get_paddle_ocr()
    if ocr is None:
        return ""

    temp_path = None
    try:
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        temp_path = tmp.name
        tmp.close()

        save_img = img.convert("RGB") if img.mode != "RGB" else img
        save_img.save(temp_path, format="PNG")

        try:
            result = ocr.ocr(temp_path, cls=True)
        except TypeError:
            result = ocr.ocr(temp_path)

        lines = []
        if result:
            for page in result:
                if not page:
                    continue
                for line in page:
                    try:
                        text = line[1][0]
                        score = 1.0
                        try:
                            score = float(line[1][1])
                        except Exception:
                            pass
                        if text and score >= 0.25:
                            lines.append(str(text))
                    except Exception:
                        continue

        return normalize_ocr_text("\n".join(lines))

    except Exception:
        return ""

    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception:
                pass


def ocr_with_tesseract(img: Image.Image) -> str:
    if pytesseract is None:
        return ""

    configs = [
        "--oem 3 --psm 6",
        "--oem 3 --psm 11",
        "--oem 3 --psm 12",
        "--oem 3 --psm 4",
    ]
    langs = ["kor+eng", "kor", "eng", None]

    best_text = ""
    for lang in langs:
        for config in configs:
            try:
                if lang:
                    text = pytesseract.image_to_string(img, lang=lang, config=config)
                else:
                    text = pytesseract.image_to_string(img, config=config)
                text = normalize_ocr_text(text)
                if len(normalize_for_match(text)) > len(normalize_for_match(best_text)):
                    best_text = text
            except Exception:
                continue

    return normalize_ocr_text(best_text)


def get_ocr_regions(img: Image.Image, top_ratio: float = 0.65) -> List[Tuple[str, Image.Image]]:
    """전체/상단/좌측/중앙 등 여러 영역을 읽어 업체명·25대고위험 누락을 줄임."""
    width, height = img.size
    top_h = max(1, int(height * top_ratio))
    mid_y1 = int(height * 0.18)
    mid_y2 = int(height * 0.78)

    regions = [
        ("full", img),
        ("top", img.crop((0, 0, width, top_h))),
        ("top_left", img.crop((0, 0, int(width * 0.72), top_h))),
        ("top_right", img.crop((int(width * 0.28), 0, width, top_h))),
        ("left", img.crop((0, 0, int(width * 0.72), height))),
        ("center", img.crop((int(width * 0.10), mid_y1, int(width * 0.90), mid_y2))),
    ]
    return regions


def append_unique_text(texts: List[str], seen: set, text: str):
    text = normalize_ocr_text(text)
    key = normalize_for_match(text)
    if text and key and key not in seen:
        texts.append(text)
        seen.add(key)


def extract_ocr_text(image_path: str, top_ratio: float = 0.65) -> str:
    """
    정확도 우선 강화 OCR.
    - 원본 이미지 기준 OCR
    - 전체/상단/좌측/중앙 영역 OCR
    - 영역별 전처리 5~7종 OCR
    - PaddleOCR 우선 + Tesseract 보조
    - 결과를 합쳐 분류 판단에 사용
    """
    try:
        img = Image.open(image_path)
        try:
            img.seek(0)
        except Exception:
            pass

        try:
            img = ImageOps.exif_transpose(img)
        except Exception:
            pass

        if img.mode != "RGB":
            img = img.convert("RGB")

        all_texts = []
        seen = set()

        paddle_available = get_paddle_ocr() is not None

        for region_name, region in get_ocr_regions(img, top_ratio=top_ratio):
            variants = make_ocr_preprocess_variants(region)

            for variant_idx, variant in enumerate(variants):
                # PaddleOCR은 정확도가 좋으므로 모든 주요 전처리 버전에 실행
                if paddle_available:
                    append_unique_text(all_texts, seen, ocr_with_paddle(variant))

                # Tesseract는 느리고 중복이 많으므로 원본/강화/이진화 일부만 보조 실행
                if variant_idx in (1, 2, 4):
                    append_unique_text(all_texts, seen, ocr_with_tesseract(variant))

        return normalize_ocr_text("\n".join(all_texts))

    except Exception:
        return ""


def extract_top_ocr_text(image_path: str, top_ratio: float = 0.65) -> str:
    """기존 함수명 호환용. 내부는 강화 OCR 사용."""
    return extract_ocr_text(image_path, top_ratio=top_ratio)


def fuzzy_contains(text: str, candidates: List[str], threshold: float = 0.72) -> bool:
    compact = normalize_for_match(text)
    if not compact:
        return False

    for cand in candidates:
        cand_norm = normalize_for_match(cand)
        if not cand_norm:
            continue
        if cand_norm in compact:
            return True

        # 긴 OCR 문자열 안에서 후보 길이만큼 잘라 유사도 검사
        n = len(cand_norm)
        if n <= 1:
            continue
        for i in range(0, max(1, len(compact) - n + 1)):
            part = compact[i:i + n]
            if SequenceMatcher(None, cand_norm, part).ratio() >= threshold:
                return True

    return False


def detect_high_risk(text: str) -> bool:
    compact = normalize_for_match(text)
    loose = str(text or "")

    if "고위험" in compact or "고 위험" in loose:
        return True
    if "25대" in compact or "25 대" in loose:
        return True

    # OCR에서 숫자/한글이 일부 깨지는 경우 보정
    high_risk_candidates = HIGH_RISK_KEYWORDS + [
        "이십오대고위험",
        "25고위험",
        "25대위험",
        "고위헙",
        "고위힘",
        "고위혐",
        "고위험작업",
        "고위험 작업",
    ]

    if fuzzy_contains(compact, high_risk_candidates, threshold=0.70):
        return True

    # 25와 위험류 단어가 따로 읽힌 경우
    has_25 = bool(re.search(r"2\s*5|25|이십오", compact))
    has_risk = fuzzy_contains(compact, ["고위험", "위험", "위헙", "위혐"], threshold=0.68)
    return has_25 and has_risk


def detect_company(text: str) -> str:
    compact = normalize_for_match(text)

    # 자주 틀리는 OCR 후보 추가
    extra_alias = {
        "원영건업": ["원영건업", "원영", "원영건", "원명건업"],
        "청암기업": ["청암기업", "청암", "청암기엽", "청암업"],
        "유셀네트웍스": ["유셀네트웍스", "유셀네트윅스", "유셀네트", "유셀", "유셀네트워크", "유셀네트웍"],
        "엠케이지": ["엠케이지", "MKG", "mkg", "엠케이", "엠케", "MK G"],
        "KEC": ["KEC", "kec", "케이이씨", "케이씨", "케이", "K E C"],
        "우신에이스": ["우신에이스", "우신", "우신에이", "우신이스"],
        "진솔": ["진솔", "진술", "진슬", "진솔건", "진솔건설"],
        "장한건설": ["장한건설", "장한", "장한전설", "장한건", "장한건썰"],
    }

    for company in COMPANY_ORDER:
        aliases = extra_alias.get(company, COMPANY_ALIAS.get(company, [company]))
        for alias in aliases:
            if normalize_for_match(alias) in compact:
                return company

    # OCR 오인식 보정: 토큰 및 슬라이딩 유사도 기반 판단
    tokens = re.findall(r"[가-힣A-Za-z0-9]{2,}", str(text or ""))
    compact_tokens = [normalize_for_match(t) for t in tokens]
    compact_tokens.append(compact)

    best_company = "기타업체"
    best_score = 0.0

    for company in COMPANY_ORDER:
        aliases = extra_alias.get(company, COMPANY_ALIAS.get(company, [company]))
        for alias in aliases:
            alias_norm = normalize_for_match(alias)
            if not alias_norm:
                continue
            for token in compact_tokens:
                if len(token) < 2:
                    continue

                # 토큰 전체 비교
                score = SequenceMatcher(None, alias_norm, token).ratio()
                if score > best_score:
                    best_score = score
                    best_company = company

                # 긴 토큰 안에 업체명이 섞여 있는 경우 부분 비교
                n = len(alias_norm)
                if len(token) >= n >= 2:
                    for i in range(0, len(token) - n + 1):
                        part = token[i:i + n]
                        score = SequenceMatcher(None, alias_norm, part).ratio()
                        if score > best_score:
                            best_score = score
                            best_company = company

    if best_score >= 0.68:
        return best_company

    return "기타업체"


def extract_sort_number(text: str) -> int:
    raw = str(text or "")

    # 엠케이지 1. 지게차 / 진솔-2 / KEC 지하 2층 등 대부분 대응
    patterns = [
        r"(?:지하|B)\s*[-]?\s*(\d+)\s*층",
        r"[-–_]\s*(\d+)",
        r"\b(\d+)\s*[.)]",
        r"(?:^|\s)(\d+)(?:\s|$)",
    ]

    for pattern in patterns:
        m = re.search(pattern, raw, flags=re.IGNORECASE)
        if m:
            try:
                return int(m.group(1))
            except Exception:
                pass

    return 0



DAILY_CLASSIFY_BATCH_SIZE = 8

DAILY_CLASSIFY_PROMPT = """
다음 이미지들은 건설현장 일일안전회의용 사진이다.
각 이미지를 업로드 순서대로 판독해서 JSON 배열만 반환하라.

각 항목:
- work_type: "material" 또는 "high_risk"
- company: 아래 업체 중 하나. 불확실하면 "기타업체"
- number: 제목/상단 표기에서 순번이 확인되면 정수, 없으면 0

분류 기준:
- "25대 고위험", "고위험", "선정사유" 등이 확인되면 high_risk
- 그렇지 않으면 material

업체 허용값:
원영건업, 청암기업, 유셀네트워크, 엠케이지, 금성, 웰시스템, KEC, 청오,
우신, MS건설, 진솔, 장한건설, 신영기초개발, KCC, 씨즌텍, 기타업체

업체명 보정:
- MKG, mkg, 엠케이, 엠케이지 -> 엠케이지
- KEC, kec, 케이이씨 -> KEC
- 유셀네트웍스, 유셀네트워크, 유셀 -> 유셀네트워크
- MS, 엠에스, 엠에스건설 -> MS건설
- 우신에이스, 우신 -> 우신

순번 표기 예:
- 업체명-1, 업체명-2
- 1. 업체명, 2. 업체명
- 업체명 1, 업체명 2

출력 예:
[
  {"work_type":"material","company":"원영건업","number":1},
  {"work_type":"high_risk","company":"엠케이지","number":2}
]

설명/코드블록 금지. 입력 이미지 개수와 출력 배열 개수는 반드시 같아야 한다.
"""


def _daily_image_data_url(uploaded_file, max_dim: int = 1500, quality: int = 82) -> str:
    uploaded_file.seek(0)
    img = Image.open(uploaded_file)
    try:
        img = ImageOps.exif_transpose(img)
    except Exception:
        pass
    if img.mode != "RGB":
        img = img.convert("RGB")

    w, h = img.size
    longest = max(w, h)
    if longest > max_dim:
        ratio = max_dim / longest
        img = img.resize(
            (max(1, int(w * ratio)), max(1, int(h * ratio))),
            Image.LANCZOS
        )

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=quality, optimize=True)
    b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
    uploaded_file.seek(0)
    return f"data:image/jpeg;base64,{b64}"


def classify_material_files_with_gpt(api_key: str, material_files) -> List[dict]:
    results = []

    for start in range(0, len(material_files), DAILY_CLASSIFY_BATCH_SIZE):
        batch = material_files[start:start + DAILY_CLASSIFY_BATCH_SIZE]

        content = [{"type": "input_text", "text": DAILY_CLASSIFY_PROMPT}]
        for i, f in enumerate(batch, start=1):
            content.append({
                "type": "input_text",
                "text": f"이미지 {i} / 파일명: {f.name}"
            })
            content.append({
                "type": "input_image",
                "image_url": _daily_image_data_url(f)
            })

        headers = {
            "Authorization": f"Bearer {api_key.strip()}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": "gpt-4o-mini",
            "input": [{"role": "user", "content": content}],
        }

        resp = requests.post(
            "https://api.openai.com/v1/responses",
            headers=headers,
            json=payload,
            timeout=75,
        )
        if resp.status_code != 200:
            raise Exception(f"일일안전회의 이미지 분류 API 오류: {resp.text}")

        data = resp.json()
        raw = _extract_openai_output_text(data)
        cleaned = str(raw or "").replace("```json", "").replace("```", "").strip()
        try:
            parsed = json.loads(cleaned)
        except Exception:
            match = re.search(r"\[.*\]", cleaned, re.S)
            parsed = json.loads(match.group(0)) if match else []

        if not isinstance(parsed, list) or len(parsed) != len(batch):
            raise ValueError(
                f"이미지 분류 개수 불일치: 입력 {len(batch)} / "
                f"출력 {len(parsed) if isinstance(parsed, list) else 0}"
            )

        for item in parsed:
            if not isinstance(item, dict):
                item = {}

            work_type = str(item.get("work_type", "material")).strip().lower()
            if work_type not in ("material", "high_risk"):
                work_type = "material"

            company = str(item.get("company", "기타업체")).strip() or "기타업체"
            company = {
                "유셀네트워크": "유셀네트웍스",
                "우신": "우신에이스",
            }.get(company, company)

            try:
                number = int(item.get("number", 0) or 0)
            except Exception:
                number = 0

            results.append({
                "work_type": work_type,
                "company": company,
                "number": number,
            })

    return results


def classify_material_work_image(
    image_path: str,
    original_name: str,
    upload_index: int,
    ocr_image_path: str = None
) -> MaterialWorkItem:
    # OCR은 원본 이미지로, PPT 삽입은 변환/축소된 JPG로 분리
    ocr_source = ocr_image_path or image_path
    ocr_text = extract_ocr_text(ocr_source, top_ratio=0.60)
    combined_text = f"{ocr_text} {original_name}"

    work_type = "high_risk" if detect_high_risk(combined_text) else "material"
    company = detect_company(combined_text)
    number = extract_sort_number(combined_text)

    return MaterialWorkItem(
        image_path=image_path,
        original_name=original_name,
        upload_index=upload_index,
        ocr_text=ocr_text,
        work_type=work_type,
        company=company,
        number=number,
    )


def company_order_index(company: str) -> int:
    try:
        return COMPANY_ORDER.index(company)
    except ValueError:
        return 999


def sort_material_work_items(items: List[MaterialWorkItem]) -> List[MaterialWorkItem]:
    # 자재입고현황 전체 → 25대 고위험작업 전체
    # 각 그룹 내부는 업체순 → 숫자순 → 업로드순
    type_order = {"material": 0, "high_risk": 1}

    return sorted(
        items,
        key=lambda x: (
            type_order.get(x.work_type, 99),
            company_order_index(x.company),
            x.number,
            x.upload_index,
        )
    )


def slide_has_text(slide, target_text: str) -> bool:
    target = normalize_text(target_text)

    for shape in iter_all_shapes(slide.shapes):
        if has_text(shape):
            if target in normalize_text(shape.text):
                return True

    for shape in iter_all_shapes(slide.shapes):
        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if target in normalize_text(cell.text):
                        return True

    return False


def find_text_target(slide, target_text: str):
    """
    PPT 플레이스홀더 탐색 강화 버전.
    1) 도형 안 텍스트가 target_text와 정확히 일치하면 반환
    2) 도형 이름(shape.name)이 target_text와 정확히 일치하면 반환
    3) 그룹도형 내부 도형까지 탐색
    4) 표 셀도 탐색

    PowerPoint에서 텍스트박스를 복사/수정하면
    화면에는 1,2,3,4가 보여도 shape.name은 TextBox 7처럼 바뀔 수 있어서
    텍스트와 도형 이름을 둘 다 확인한다.
    """
    target = normalize_text(target_text)
    target_match = normalize_for_match(target_text)

    # 1차: 텍스트 정확 일치
    for shape in iter_all_shapes(slide.shapes):
        if has_text(shape):
            if normalize_text(shape.text) == target:
                return ("shape", shape)

    # 2차: 도형 이름 정확 일치
    for shape in iter_all_shapes(slide.shapes):
        shape_name = normalize_text(getattr(shape, "name", ""))
        if shape_name == target:
            return ("shape", shape)

    # 3차: 공백/줄바꿈 제거 후 텍스트 일치
    for shape in iter_all_shapes(slide.shapes):
        if has_text(shape):
            if normalize_for_match(shape.text) == target_match:
                return ("shape", shape)

    # 4차: 공백/줄바꿈 제거 후 도형 이름 일치
    for shape in iter_all_shapes(slide.shapes):
        shape_name = getattr(shape, "name", "")
        if normalize_for_match(shape_name) == target_match:
            return ("shape", shape)

    # 5차: 표 셀 탐색
    for shape in iter_all_shapes(slide.shapes):
        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if normalize_text(cell.text) == target:
                        return ("cell", cell)
                    if normalize_for_match(cell.text) == target_match:
                        return ("cell", cell)

    return None


def debug_slide_placeholders(slide):
    """Streamlit 화면에서 PPT 도형 이름/텍스트 확인용."""
    rows = []
    for idx, shape in enumerate(iter_all_shapes(slide.shapes), start=1):
        rows.append({
            "순번": idx,
            "도형이름": getattr(shape, "name", ""),
            "텍스트": normalize_text(shape.text) if has_text(shape) else "",
            "타입": str(getattr(shape, "shape_type", "")),
        })
    return rows



def find_tbm_translation_cells(slide):
    """TBM 템플릿 표 구조용 안전장치.
    sample_template.pptx는 1,2,3,4가 표 셀에 들어있으므로,
    텍스트 탐색이 실패해도 표의 우측 2~5행 셀을 직접 사용한다.
    """
    for shape in iter_all_shapes(slide.shapes):
        try:
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                if len(table.rows) >= 6 and len(table.columns) >= 2:
                    # 기본 템플릿 구조: 2~5행, 우측 열이 1~4 번역칸
                    return {
                        "1": ("cell", table.cell(2, 1)),
                        "2": ("cell", table.cell(3, 1)),
                        "3": ("cell", table.cell(4, 1)),
                        "4": ("cell", table.cell(5, 1)),
                    }
        except Exception:
            continue
    return {}


def find_slide_index_by_text(prs, target_text: str):
    for idx, slide in enumerate(prs.slides):
        if slide_has_text(slide, target_text):
            return idx
    return None


def set_target_text(
    target_obj,
    text: str,
    size_pt: int,
    font_name: str = None,
    bold: bool = False,
    font_color=None
):
    kind, obj = target_obj

    tf = obj.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold

    if font_name:
        run.font.name = font_name

    if font_color:
        run.font.color.rgb = font_color


def add_picture_to_shape(slide, image_path, target_shape):
    """
    PHOTO_BOX에 사진을 넣을 때 회전/왜곡 금지.
    - 사진 비율 유지
    - 박스는 꽉 채움(cover crop)
    - crop만 적용하고 rotation은 항상 0
    """
    try:
        img = Image.open(image_path)
        try:
            img = ImageOps.exif_transpose(img)
        except Exception:
            pass
        img_w, img_h = img.size
        img.close()
    except Exception:
        slide.shapes.add_picture(
            image_path,
            target_shape.left,
            target_shape.top,
            width=target_shape.width,
            height=target_shape.height
        )
        return

    if img_w <= 0 or img_h <= 0:
        return

    box_w = float(target_shape.width)
    box_h = float(target_shape.height)
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h

    pic = slide.shapes.add_picture(
        image_path,
        target_shape.left,
        target_shape.top,
        width=target_shape.width,
        height=target_shape.height
    )

    try:
        pic.rotation = 0
    except Exception:
        pass

    # python-pptx crop 값은 0~1 비율. 가로/세로를 바꾸지 않고 초과분만 잘라낸다.
    try:
        if img_ratio > box_ratio:
            crop = max(0.0, min(0.49, (1 - (box_ratio / img_ratio)) / 2))
            pic.crop_left = crop
            pic.crop_right = crop
            pic.crop_top = 0
            pic.crop_bottom = 0
        elif img_ratio < box_ratio:
            crop = max(0.0, min(0.49, (1 - (img_ratio / box_ratio)) / 2))
            pic.crop_top = crop
            pic.crop_bottom = crop
            pic.crop_left = 0
            pic.crop_right = 0
    except Exception:
        pass

    return pic


def duplicate_slide(prs, source_slide):
    """
    슬라이드 복제 시 도형 XML뿐 아니라 이미지 relationship도 같이 복사.
    기존 deepcopy만 쓰면 국기 같은 그림이 2번째 슬라이드부터 깨질 수 있다.
    """
    blank_slide_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_slide_layout)

    for shape in source_slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(
            new_el,
            "p:extLst"
        )

    # 이미지/차트 등 외부 관계 복사. notesSlide는 복사하지 않음.
    try:
        for rel in source_slide.part.rels.values():
            if "notesSlide" in rel.reltype:
                continue
            try:
                new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
            except Exception:
                try:
                    new_slide.part.rels.add_relationship(rel.reltype, rel.target_part, rel.rId)
                except Exception:
                    pass
    except Exception:
        pass

    return new_slide


def get_slide_index_by_id(prs, slide_id: int):
    for idx, slide in enumerate(prs.slides):
        if slide.slide_id == slide_id:
            return idx
    return None


def delete_slide_by_index(prs, idx: int):
    slide_id = prs.slides._sldIdLst[idx]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[idx]


def delete_slide_by_id(prs, slide_id: int):
    idx = get_slide_index_by_id(prs, slide_id)
    if idx is not None:
        delete_slide_by_index(prs, idx)


def move_slide_by_id(prs, slide_id: int, target_index: int):
    idx = get_slide_index_by_id(prs, slide_id)
    if idx is None:
        return

    sld_id = prs.slides._sldIdLst[idx]
    prs.slides._sldIdLst.remove(sld_id)

    target_index = max(0, min(target_index, len(prs.slides._sldIdLst)))
    prs.slides._sldIdLst.insert(target_index, sld_id)


def find_slide_index_by_any_text(prs, texts: List[str]):
    for text in texts:
        idx = find_slide_index_by_text(prs, text)
        if idx is not None:
            return idx
    return None


def clone_and_fill_bad_slides(prs, template_slide, bad_items: List[DailySlideData]) -> List[int]:
    created_ids = []
    for item in bad_items:
        new_slide = duplicate_slide(prs, template_slide)
        fill_daily_slide(new_slide, item, strict=False)
        created_ids.append(new_slide.slide_id)
    return created_ids


def clone_and_fill_material_slides(prs, template_slide, material_items: List[MaterialWorkItem]) -> List[int]:
    created_ids = []
    for item in material_items:
        new_slide = duplicate_slide(prs, template_slide)
        insert_image_to_placeholder(new_slide, TIME_BOX_TEXT, item.image_path)
        created_ids.append(new_slide.slide_id)
    return created_ids


def fill_slide_by_placeholders(slide, item: SlideData, strict: bool = True):
    photo_target = find_text_target(slide, PHOTO_BOX_TEXT)
    ko_target = find_text_target(slide, KO_BOX_TEXT)
    zh_target = find_text_target(slide, ZH_BOX_TEXT)
    vi_target = find_text_target(slide, VI_BOX_TEXT)
    my_target = find_text_target(slide, MY_BOX_TEXT)

    # 안전장치: sample_template.pptx의 1,2,3,4가 표 셀로 들어간 경우 직접 매칭
    table_cells = find_tbm_translation_cells(slide)
    if ko_target is None:
        ko_target = table_cells.get("1")
    if zh_target is None:
        zh_target = table_cells.get("2")
    if vi_target is None:
        vi_target = table_cells.get("3")
    if my_target is None:
        my_target = table_cells.get("4")

    missing = []
    for name, obj in [
        ("PHOTO_BOX", photo_target),
        ("1", ko_target),
        ("2", zh_target),
        ("3", vi_target),
        ("4", my_target),
    ]:
        if obj is None:
            missing.append(name)

    if missing:
        if strict:
            debug = []
            for shape in iter_all_shapes(slide.shapes):
                try:
                    debug.append(f"name=[{getattr(shape, 'name', '')}] text=[{getattr(shape, 'text', '')}]")
                    if hasattr(shape, "has_table") and shape.has_table:
                        for r, row in enumerate(shape.table.rows):
                            for c, cell in enumerate(row.cells):
                                debug.append(f"table[{r},{c}]=[{cell.text}]")
                except Exception:
                    pass
            raise ValueError(
                "슬라이드에서 플레이스홀더를 찾지 못했습니다: "
                + ", ".join(missing)
                + "\n\n"
                + "\n".join(debug[:80])
            )
        return

    photo_kind, photo_obj = photo_target
    if photo_kind != "shape":
        if strict:
            raise ValueError("PHOTO_BOX는 텍스트 상자/도형이어야 합니다.")
        return

    add_picture_to_shape(slide, item.image_path, photo_obj)

    set_target_text(ko_target, item.ko, BASE_FONT_SIZE_PT)
    set_target_text(zh_target, item.zh, BASE_FONT_SIZE_PT)
    set_target_text(vi_target, item.vi, BASE_FONT_SIZE_PT)
    set_target_text(my_target, item.my, BASE_FONT_SIZE_PT)


def fill_daily_slide(slide, item: DailySlideData, strict: bool = False):
    photo_target = find_text_target(slide, DAILY_PHOTO_BOX_TEXT)
    text_target = find_text_target(slide, DAILY_TEXT_BOX_TEXT)

    if photo_target:
        kind, obj = photo_target
        if kind == "shape":
            add_picture_to_shape(slide, item.image_path, obj)
    elif strict:
        raise ValueError("PHOTO_BOX_1 플레이스홀더를 찾지 못했습니다.")

    if text_target:
        set_target_text(
            text_target,
            item.text,
            28,
            font_name="맑은 고딕",
            bold=True,
            font_color=RGBColor(0, 0, 0)
        )
    elif strict:
        raise ValueError("TEXT_BOX_1 플레이스홀더를 찾지 못했습니다.")


def insert_image_to_placeholder(slide, placeholder_text: str, image_path: str):
    target = find_text_target(slide, placeholder_text)

    if target is None:
        return

    kind, obj = target

    if kind != "shape":
        return

    add_picture_to_shape(slide, image_path, obj)


def fill_material_slides(prs, material_items: List[MaterialWorkItem]):
    if not material_items:
        return

    base_idx = find_slide_index_by_text(prs, TIME_BOX_TEXT)

    if base_idx is None:
        return

    base_slide = prs.slides[base_idx]

    for i, item in enumerate(material_items):
        if i == 0:
            target_slide = base_slide
        else:
            target_slide = duplicate_slide(prs, base_slide)

        insert_image_to_placeholder(
            target_slide,
            TIME_BOX_TEXT,
            item.image_path
        )


def fill_date_box(slide):
    target = find_text_target(slide, DATE_BOX_TEXT)

    if target:
        set_target_text(
            target,
            get_korean_date_text(),
            30,
            font_name="맑은 고딕",
            bold=True,
            font_color=RGBColor(0, 0, 0)
        )


def delete_extra_slides(prs, keep_slide_count: int):
    for idx in range(len(prs.slides) - 1, keep_slide_count - 1, -1):
        slide = prs.slides[idx]

        if slide_has_text(slide, HOLD_POINT_TEXT):
            continue

        slide_id = prs.slides._sldIdLst[idx]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[idx]


def build_ppt(slide_data_list: List[SlideData]) -> io.BytesIO:
    if not os.path.exists(TEMPLATE_PPT):
        raise FileNotFoundError(f"템플릿 파일이 없습니다: {TEMPLATE_PPT}")

    prs = Presentation(TEMPLATE_PPT)

    if not slide_data_list:
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return out

    if len(prs.slides) < 1:
        raise ValueError("TBM 템플릿에 기준 슬라이드가 없습니다.")

    # 국기 이미지 깨짐 방지: 템플릿에 이미 들어있는 슬라이드를 먼저 사용한다.
    # 템플릿 수를 초과한 사진만 relationship 복사 방식으로 추가 복제한다.
    filled_ids = []
    template_slide_count = len(prs.slides)
    base_slide = prs.slides[0]

    for i, item in enumerate(slide_data_list):
        if i < template_slide_count:
            target_slide = prs.slides[i]
        else:
            target_slide = duplicate_slide(prs, base_slide)

        fill_slide_by_placeholders(target_slide, item, strict=True)
        filled_ids.append(target_slide.slide_id)

        if i % 10 == 0:
            gc.collect()

    keep_ids = set(filled_ids)
    for idx in range(len(prs.slides) - 1, -1, -1):
        slide = prs.slides[idx]
        if slide.slide_id in keep_ids:
            continue
        if slide_has_text(slide, HOLD_POINT_TEXT):
            continue
        delete_slide_by_index(prs, idx)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    gc.collect()
    return out


def build_daily_ppt(
    bad_items: List[DailySlideData],
    material_items: List[MaterialWorkItem]
) -> io.BytesIO:
    if not os.path.exists(DAILY_TEMPLATE_PPT):
        raise FileNotFoundError(f"템플릿 파일이 없습니다: {DAILY_TEMPLATE_PPT}")

    prs = Presentation(DAILY_TEMPLATE_PPT)

    if len(prs.slides) >= 1:
        fill_date_box(prs.slides[0])

    # 네이버 날씨 자동 캡처 기능은 제거됨.
    # 템플릿의 날씨 관련 슬라이드/영역은 원본 상태 그대로 유지한다.

    # 템플릿 기준 슬라이드 찾기
    bad_template_idx = find_slide_index_by_text(prs, DAILY_PHOTO_BOX_TEXT)
    material_template_idx = find_slide_index_by_text(prs, TIME_BOX_TEXT)

    if bad_items and bad_template_idx is None:
        raise ValueError("부적합사진 기준 슬라이드(PHOTO_BOX_1)를 찾지 못했습니다.")

    if material_items and material_template_idx is None:
        raise ValueError("자재입고/고위험 기준 슬라이드(TIME_BOX_1)를 찾지 못했습니다.")

    bad_template_slide = prs.slides[bad_template_idx] if bad_template_idx is not None else None
    material_template_slide = prs.slides[material_template_idx] if material_template_idx is not None else None

    bad_template_id = bad_template_slide.slide_id if bad_template_slide is not None else None
    material_template_id = material_template_slide.slide_id if material_template_slide is not None else None

    created_bad_ids = []
    created_material_ids = []

    if bad_template_slide is not None:
        created_bad_ids = clone_and_fill_bad_slides(prs, bad_template_slide, bad_items)

    if material_template_slide is not None:
        created_material_ids = clone_and_fill_material_slides(prs, material_template_slide, material_items)

    # 원본 기준 슬라이드는 빈 템플릿이므로 제거. 같은 슬라이드 중복 제거 방지.
    for sid in sorted({x for x in [bad_template_id, material_template_id] if x is not None}, reverse=True):
        delete_slide_by_id(prs, sid)

    # 동적 슬라이드를 명일 작업내용 발표 앞에 배치.
    # 명일 문구를 못 찾으면 HOLD POINT 앞, 그것도 못 찾으면 맨 뒤에 배치.
    anchor_idx = find_slide_index_by_any_text(
        prs,
        [
            "명일 작업내용 발표",
            "명일작업내용발표",
            "명일 작업내용",
            "명일작업내용",
        ]
    )

    if anchor_idx is None:
        anchor_idx = find_slide_index_by_text(prs, HOLD_POINT_TEXT)

    if anchor_idx is None:
        anchor_idx = len(prs.slides)

    desired_dynamic_ids = created_bad_ids + created_material_ids

    insert_at = anchor_idx
    for sid in desired_dynamic_ids:
        move_slide_by_id(prs, sid, insert_at)
        insert_at += 1

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)


    return out


def render_tbm_input_area():
    files = st.file_uploader(
        "사진 업로드",
        accept_multiple_files=True,
        type=["jpg", "png", "jpeg", "webp", "heic", "heif", "mpo"],
        key="main_tbm_uploader"
    )

    if files:
        if len(files) > MAX_TBM_FILES_SOFT_WARN:
            st.warning(f"사진이 {len(files)}장입니다. 고용량 사진은 자동 축소해서 처리하지만 생성 시간이 길어질 수 있습니다.")

        slide_inputs = []
        temp_paths = []

        for idx, f in enumerate(files):
            with st.expander(f"슬라이드 #{idx+1}", expanded=True):
                c1, c2 = st.columns([1, 4])

                suffix = os.path.splitext(f.name)[1].lower() or ".jpg"

                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                    tmp.write(f.getbuffer())
                    original_path = tmp.name
                    temp_paths.append(original_path)

                jpg_path = convert_to_jpg(original_path, max_size=TBM_IMAGE_MAX_SIZE, quality=TBM_IMAGE_QUALITY)
                temp_paths.append(jpg_path)

                c1.image(jpg_path, width=150)

                ko_input = c2.text_input(
                    "한국어 문구",
                    value="",
                    placeholder="예: 지정된 이동통로 통행",
                    key=f"main_tbm_ko_{idx}"
                )

                slide_inputs.append(SlideData(jpg_path, ko_input, "", "", "", get_image_orientation(jpg_path)))

        if st.button("PPT 생성", key="main_create_btn"):
            try:
                if "GPT_API_KEY" not in st.secrets:
                    raise ValueError("Secrets에 GPT_API_KEY 설정 필요")

                with st.spinner("번역 중..."):
                    ko_list = [s.ko for s in slide_inputs]

                    if any(not x.strip() for x in ko_list):
                        raise ValueError("빈 한국어 문구가 있습니다. 모든 슬라이드 문구를 입력하세요.")

                    translations = translate_all_with_gpt(
                        st.secrets["GPT_API_KEY"],
                        ko_list
                    )

                    for s, tr in zip(slide_inputs, translations):
                        s.zh = tr["zh"]
                        s.vi = tr["vi"]
                        s.my = tr["my"]

                with st.spinner("PPT 생성 중..."):
                    ppt = build_ppt(slide_inputs)

                save_generated_ppt_to_bad_photo_storage(ppt, OUTPUT_PPT_NAME)

                st.success("완료!")
                st.download_button(
                    "PPT 다운로드",
                    ppt,
                    file_name=OUTPUT_PPT_NAME,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="main_download_btn"
                )

            except Exception as e:
                st.error(f"오류 발생: {e}")

            finally:
                for p in temp_paths:
                    if os.path.exists(p):
                        try:
                            os.remove(p)
                        except Exception:
                            pass


def render_daily_safety_meeting():
    st.markdown(
        """
        <h2 style="color:#d00000; font-weight:800; font-size:48px; margin-top:0.35rem; margin-bottom:0.35rem;">
            일일안전회의
        </h2>
        """,
        unsafe_allow_html=True
    )

    # ------------------------------------------------------------
    # 중요:
    # 예전 코드는 자재입고 사진을 올리는 즉시 모든 사진에 강화 OCR을 실행한 뒤
    # st.button()까지 내려갔다. 그래서 화면에는 버튼이 보이더라도 Streamlit 스크립트가
    # 계속 실행 중이라 버튼이 회색/비활성처럼 보이고 실제로 클릭할 수 없었다.
    #
    # 이제 업로드 단계에서는 OCR/고용량 변환을 절대 하지 않는다.
    # 버튼을 먼저 즉시 렌더링하고, 사용자가 버튼을 누른 뒤에만 OCR + PPT 생성을 실행한다.
    # ------------------------------------------------------------

    bad_files = st.file_uploader(
        "부적합사진",
        accept_multiple_files=True,
        type=["jpg", "png", "jpeg", "webp", "heic", "heif", "mpo"],
        key="daily_bad_uploader"
    )

    bad_text_values = []

    if bad_files:
        if len(bad_files) > MAX_DAILY_FILES_SOFT_WARN:
            st.warning(
                f"부적합사진이 {len(bad_files)}장입니다. "
                "생성 버튼을 누른 뒤 자동 축소 처리합니다."
            )

        st.markdown("#### 부적합사진")

        for idx, f in enumerate(bad_files):
            with st.expander(f"부적합사진 #{idx + 1}", expanded=False):
                c1, c2 = st.columns([1, 4])

                with c1:
                    # 업로드 단계에서는 임시 JPG 변환 없이 원본 미리보기만 표시.
                    try:
                        st.image(f, width=130)
                    except Exception:
                        st.caption(f.name)

                with c2:
                    text_value = st.text_input(
                        "문구 입력",
                        value="",
                        placeholder="예: 자재 반입 확인",
                        key=f"daily_bad_text_{idx}"
                    )
                    st.caption(f.name)

            bad_text_values.append(text_value)

    material_files = st.file_uploader(
        "자재입고 및 고위험작업",
        accept_multiple_files=True,
        type=["jpg", "png", "jpeg", "webp", "heic", "heif", "mpo"],
        key="daily_material_uploader"
    )

    if material_files:
        if len(material_files) > MAX_DAILY_FILES_SOFT_WARN:
            st.warning(
                f"자재입고 및 고위험작업 사진이 {len(material_files)}장입니다. "
                "OCR은 생성 버튼을 누른 뒤 실행합니다."
            )

        st.markdown("#### 자재입고 및 고위험작업")

        # 여기서는 OCR을 하지 않는다. 파일명/간단 미리보기만 보여준다.
        for idx, f in enumerate(material_files):
            c1, c2 = st.columns([1, 4])
            with c1:
                try:
                    st.image(f, width=110)
                except Exception:
                    st.caption("사진")
            with c2:
                st.caption(f"{idx + 1}. {f.name}")
                st.caption("분류·업체명·순서는 PPT 생성 시 OCR로 자동 판독")

    # 업로더 바로 다음에 버튼을 즉시 렌더링한다.
    st.markdown("<div style='height:0.35rem'></div>", unsafe_allow_html=True)

    daily_create_clicked = st.button(
        "일일안전회의 PPT 생성",
        key="daily_create_btn",
        use_container_width=True,
        type="primary"
    )

    st.markdown("<div style='height:0.6rem'></div>", unsafe_allow_html=True)

    if daily_create_clicked:
        bad_items = []
        material_items = []
        temp_paths = []

        try:
            # 1. 부적합사진 변환
            if bad_files:
                with st.spinner("부적합사진 처리 중..."):
                    for idx, f in enumerate(bad_files):
                        suffix = os.path.splitext(f.name)[1].lower() or ".jpg"

                        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                            tmp.write(f.getvalue())
                            original_path = tmp.name
                            temp_paths.append(original_path)

                        jpg_path = convert_to_jpg(
                            original_path,
                            max_size=DAILY_IMAGE_MAX_SIZE,
                            quality=DAILY_IMAGE_QUALITY
                        )
                        temp_paths.append(jpg_path)

                        text_value = bad_text_values[idx] if idx < len(bad_text_values) else ""
                        bad_items.append(
                            DailySlideData(
                                jpg_path,
                                text_value,
                                get_image_orientation(jpg_path)
                            )
                        )

            # 2. 자재입고/25대 고위험 고속 분류
            if material_files:
                if "GPT_API_KEY" not in st.secrets:
                    raise ValueError("Secrets에 GPT_API_KEY 설정 필요")

                total = len(material_files)
                progress = st.progress(0, text="사진 분류 준비 중...")

                classification_results = []

                for batch_start in range(0, total, DAILY_CLASSIFY_BATCH_SIZE):
                    batch_end = min(batch_start + DAILY_CLASSIFY_BATCH_SIZE, total)
                    progress.progress(
                        batch_start / max(1, total),
                        text=f"사진 분류 중... {batch_start + 1}~{batch_end}/{total}"
                    )
                    batch_files = material_files[batch_start:batch_end]
                    batch_results = classify_material_files_with_gpt(
                        st.secrets["GPT_API_KEY"],
                        batch_files
                    )
                    classification_results.extend(batch_results)

                progress.progress(0.72, text="분류 완료 · PPT용 사진 변환 중...")

                for idx, (f, cls) in enumerate(zip(material_files, classification_results)):
                    suffix = os.path.splitext(f.name)[1].lower() or ".jpg"

                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                        tmp.write(f.getvalue())
                        material_original_path = tmp.name
                        temp_paths.append(material_original_path)

                    material_jpg_path = convert_to_jpg(
                        material_original_path,
                        max_size=DAILY_IMAGE_MAX_SIZE,
                        quality=DAILY_IMAGE_QUALITY
                    )
                    temp_paths.append(material_jpg_path)

                    material_items.append(
                        MaterialWorkItem(
                            image_path=material_jpg_path,
                            original_name=f.name,
                            upload_index=idx,
                            ocr_text="",
                            work_type=cls["work_type"],
                            company=cls["company"],
                            number=cls["number"],
                        )
                    )

                    progress.progress(
                        0.72 + 0.25 * ((idx + 1) / max(1, total)),
                        text=f"PPT용 사진 준비 중... {idx + 1}/{total}"
                    )

                progress.progress(1.0, text="분류 및 사진 준비 완료")

            # 3. 정렬 + PPT 생성
            with st.spinner("PPT 생성 중..."):
                sorted_material_items = sort_material_work_items(material_items)
                ppt = build_daily_ppt(bad_items, sorted_material_items)

            save_generated_ppt_to_bad_photo_storage(
                ppt,
                DAILY_OUTPUT_PPT_NAME
            )

            st.success("완료!")

            st.download_button(
                "일일안전회의 PPT 다운로드",
                ppt,
                file_name=DAILY_OUTPUT_PPT_NAME,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="daily_download_btn",
                use_container_width=True
            )

        except Exception as e:
            st.error(f"오류 발생: {e}")

        finally:
            for p in temp_paths:
                if os.path.exists(p):
                    try:
                        os.remove(p)
                    except Exception:
                        pass

    # 생성 버튼과 다음 섹션 사이 여백
    st.markdown("<div style='height:1.2rem'></div>", unsafe_allow_html=True)
    st.markdown("---")
    render_bad_photo_storage()
    st.markdown("---")

def load_shared_notice() -> str:
    try:
        if os.path.exists(SHARED_NOTICE_FILE):
            with open(SHARED_NOTICE_FILE, "r", encoding="utf-8") as f:
                return f.read()
    except Exception:
        pass
    return ""


def save_shared_notice(text: str):
    try:
        with open(SHARED_NOTICE_FILE, "w", encoding="utf-8") as f:
            f.write(text or "")

        meta = {
            "saved_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }
        with open(SHARED_NOTICE_META_FILE, "w", encoding="utf-8") as f:
            json.dump(meta, f, ensure_ascii=False, indent=2)

        return True
    except Exception:
        return False


def clear_shared_notice():
    ok = save_shared_notice("")
    return ok


def load_shared_notice_saved_at() -> str:
    try:
        if os.path.exists(SHARED_NOTICE_META_FILE):
            with open(SHARED_NOTICE_META_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            return str(data.get("saved_at", "") or "")
    except Exception:
        pass
    return ""


def render_shared_notice_board():
    st.markdown("---")
    st.markdown(
        """
        <h2 style="font-size:1.35em; font-weight:700; margin-top:0.35rem; margin-bottom:0.35rem;">
            공지사항 / 메모
        </h2>
        """,
        unsafe_allow_html=True
    )

    current_text = load_shared_notice()

    notice_text = st.text_area(
        "",
        value=current_text,
        height=260,
        key="shared_notice_text_area",
        placeholder="공지사항/메모/안전사항 등."
    )

    col_save, col_delete, col_space = st.columns([0.8, 0.8, 5.4])

    with col_save:
        if st.button("저장", use_container_width=True, key="shared_notice_save_btn"):
            if save_shared_notice(notice_text):
                st.success("저장 완료")
                st.rerun()
            else:
                st.error("저장 실패")

    with col_delete:
        if st.button("삭제", use_container_width=True, key="shared_notice_delete_btn"):
            if clear_shared_notice():
                st.success("삭제 완료")
                st.rerun()
            else:
                st.error("삭제 실패")

    saved_at = load_shared_notice_saved_at()
    if saved_at:
        st.caption(f"최종 저장: {saved_at}")
    else:
        st.caption("최종 저장: 없음")


# ============================================================
# 체감온도 측정 기록 (업로드된 관리대장 템플릿 기반, GPT Vision OCR)
# ============================================================
#
# 필요 설정:
#   1) 저장소의 templates/ 폴더에 heat_index_template.xlsx 추가
#   2) Secrets에 GPT_API_KEY 필요 (기존 번역 기능과 동일한 키 재사용)
#
# 동작 방식:
#   - 사진 1장 = 카카오톡 캡쳐 등, 측정 기록이 여러 건 섞여있을 수 있음
#     → GPT-4o-mini Vision으로 사진 1장에서 기록 여러 건을 한번에 구조화 추출
#   - 시트 1개 = 특정 "측정장소 + 날짜" 하루치 (템플릿 그대로 복제, NO 1~9 슬롯)
#   - 같은 장소/날짜는 표기가 조금 달라도 같은 시트로 병합
#   - "2. 조치사항" 섹션은 템플릿 내용을 그대로 유지 (사진과 무관하므로 손대지 않음)
#   - 사진에 체감온도가 안 보이면, 기온·습도로 기상청 공식 체감온도 산출식을 이용해 계산
#     (표를 보고 어림잡지 않음 — 실측된 기온·습도로부터 결정적으로 계산되는 값이라 지어내는 것이 아님)
#   - 측정값 판독 공백은 같은 장소·날짜의 정상값 평균으로 자동 보완하고 비고란에 표시
#     (같은 장소 자료가 없으면 같은 날짜의 다른 장소 평균을 보조로 사용)
#   - 측정이 2시간 넘게 비면 비고란에 "간격초과" 메모를 함께 남김
#   - 사진을 올리면 즉시(버튼 클릭 없이) 자동 저장 → 페이지를 나가도 작업이 유실되지 않음
#     이후 확인해서 값이 틀렸으면 "수정 반영"으로 같은 칸을 덮어씀 (중복 저장 안 됨)

HEAT_TEMPLATE_XLSX = os.path.join(BASE_DIR, "templates", "heat_index_template.xlsx")
HEAT_LOG_FILE = os.path.join(BASE_DIR, "heat_index_log.xlsx")
HEAT_UPLOAD_HISTORY_FILE = os.path.join(BASE_DIR, "heat_upload_history.json")
HEAT_EXPORT_HISTORY_DIR = os.path.join(BASE_DIR, "heat_export_history")
HEAT_LOG_GAP_MINUTES = 120
HEAT_PROCESSING_STALE_SECONDS = 15 * 60

# 같은 파일의 반복 업로드와 동시에 들어오는 중복 저장을 막기 위한 프로세스 내부 잠금.
_HEAT_HISTORY_LOCK = threading.RLock()


def calc_heat_index(Ta: float, RH: float) -> float:
    """기상자료개방포털 여름철 체감온도 공식(2022.6.2 개정판).
    습구온도는 Stull(2011) 근사식으로 추정.
    체감온도 = -0.2442 + 0.55399Tw + 0.45535Ta - 0.0022Tw^2 + 0.00278TwTa + 3.0
    """
    Tw = (Ta * math.atan(0.151977 * (RH + 8.313659) ** 0.5)
          + math.atan(Ta + RH)
          - math.atan(RH - 1.676331)
          + 0.00391838 * (RH ** 1.5) * math.atan(0.023101 * RH)
          - 4.686035)
    hi = -0.2442 + 0.55399 * Tw + 0.45535 * Ta - 0.0022 * (Tw ** 2) + 0.00278 * Tw * Ta + 3.0
    return round(hi, 1)


def pil_image_to_data_url(
    img: Image.Image,
    max_dim: int = 2600,
    min_dim: int = 1600,
    quality: int = 94,
) -> str:
    """PIL 이미지를 GPT Vision용 data URL로 변환.

    작은 LCD 사진은 숫자 세그먼트가 뭉개지지 않도록 확대하고,
    큰 이미지는 메모리/전송량 보호를 위해 축소한다.
    """
    if img.mode != "RGB":
        img = img.convert("RGB")

    width, height = img.size
    longest = max(width, height)

    if longest < min_dim:
        scale = min_dim / max(1, longest)
        img = img.resize(
            (max(1, int(width * scale)), max(1, int(height * scale))),
            Image.LANCZOS,
        )
    elif longest > max_dim:
        scale = max_dim / longest
        img = img.resize(
            (max(1, int(width * scale)), max(1, int(height * scale))),
            Image.LANCZOS,
        )

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=quality, optimize=True)
    b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
    return f"data:image/jpeg;base64,{b64}"


def image_to_data_url(image_path: str, max_dim: int = 2400, quality: int = 92) -> str:
    """GPT Vision 전송용 원본 이미지 data URL.

    이전보다 작은 사진을 적극적으로 확대해 실제 온습도계 LCD 숫자를
    바로 판독할 수 있도록 한다.
    """
    img = Image.open(image_path)
    img = ImageOps.exif_transpose(img)
    return pil_image_to_data_url(
        img,
        max_dim=max_dim,
        min_dim=1600,
        quality=quality,
    )


def _heat_rect_iou(a, b) -> float:
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ix1, iy1 = max(ax1, bx1), max(ay1, by1)
    ix2, iy2 = min(ax2, bx2), min(ay2, by2)
    iw, ih = max(0, ix2 - ix1), max(0, iy2 - iy1)
    inter = iw * ih
    if inter <= 0:
        return 0.0
    area_a = max(1, (ax2 - ax1) * (ay2 - ay1))
    area_b = max(1, (bx2 - bx1) * (by2 - by1))
    return inter / max(1, area_a + area_b - inter)


def detect_heat_lcd_candidate_crops(img: Image.Image, max_candidates: int = 3) -> List[Image.Image]:
    """사진에서 LCD로 보이는 고대비 직사각형 영역을 자동 탐색해 확대용 크롭을 반환."""
    if img.mode != "RGB":
        img = img.convert("RGB")

    original_w, original_h = img.size
    if original_w < 80 or original_h < 80:
        return []

    preview = img.copy()
    preview.thumbnail((720, 720), Image.LANCZOS)
    pw, ph = preview.size

    gray = ImageOps.autocontrast(preview.convert("L"), cutoff=1)
    edges = ImageOps.autocontrast(gray.filter(ImageFilter.FIND_EDGES), cutoff=1)

    candidates = []
    for width_ratio in (0.14, 0.19, 0.25, 0.33, 0.43):
        win_w = max(70, int(pw * width_ratio))
        for aspect in (0.9, 1.15, 1.45, 1.8):
            win_h = max(55, int(win_w / aspect))
            if win_w >= pw or win_h >= ph:
                continue

            step = max(26, int(min(win_w, win_h) * 0.42))
            y_positions = list(range(0, max(1, ph - win_h + 1), step))
            x_positions = list(range(0, max(1, pw - win_w + 1), step))
            if not y_positions or y_positions[-1] != ph - win_h:
                y_positions.append(ph - win_h)
            if not x_positions or x_positions[-1] != pw - win_w:
                x_positions.append(pw - win_w)

            for y1 in y_positions:
                for x1 in x_positions:
                    box = (x1, y1, x1 + win_w, y1 + win_h)
                    edge_stat = ImageStat.Stat(edges.crop(box))
                    gray_stat = ImageStat.Stat(gray.crop(box))
                    edge_mean = edge_stat.mean[0]
                    contrast = gray_stat.stddev[0]
                    brightness = gray_stat.mean[0]

                    cx = (x1 + win_w / 2) / max(1, pw)
                    cy = (y1 + win_h / 2) / max(1, ph)
                    center_bonus = max(0.0, 1.0 - ((cx - 0.5) ** 2 + (cy - 0.5) ** 2) ** 0.5 * 1.4)
                    brightness_penalty = 8.0 if brightness < 18 or brightness > 242 else 0.0
                    score = edge_mean * 0.72 + contrast * 0.78 + center_bonus * 7.0 - brightness_penalty
                    candidates.append((score, box))

    selected = []
    for score, box in sorted(candidates, key=lambda x: x[0], reverse=True):
        if any(_heat_rect_iou(box, prev_box) > 0.35 for _, prev_box in selected):
            continue
        selected.append((score, box))
        if len(selected) >= max_candidates:
            break

    scale_x = original_w / max(1, pw)
    scale_y = original_h / max(1, ph)
    crops = []
    for _, (x1, y1, x2, y2) in selected:
        ox1, oy1 = int(x1 * scale_x), int(y1 * scale_y)
        ox2, oy2 = int(x2 * scale_x), int(y2 * scale_y)
        pad_x = int((ox2 - ox1) * 0.18)
        pad_y = int((oy2 - oy1) * 0.22)
        ox1, oy1 = max(0, ox1 - pad_x), max(0, oy1 - pad_y)
        ox2, oy2 = min(original_w, ox2 + pad_x), min(original_h, oy2 + pad_y)
        if ox2 - ox1 >= 50 and oy2 - oy1 >= 40:
            crops.append(img.crop((ox1, oy1, ox2, oy2)))
    return crops


def make_heat_meter_vision_images(image_path: str) -> List[str]:
    """LCD 자동 탐색 영역과 원본 보정본을 생성해 흐린 실측기 숫자를 재판독."""
    img = Image.open(image_path)
    img = ImageOps.exif_transpose(img)
    if img.mode != "RGB":
        img = img.convert("RGB")

    variants = [img.copy()]

    # 전체 사진 강화본: 카카오톡 캡처의 이름·시간·위치 문맥도 유지한다.
    full_enhanced = ImageOps.autocontrast(img, cutoff=1)
    full_enhanced = ImageEnhance.Contrast(full_enhanced).enhance(1.40)
    full_enhanced = ImageEnhance.Sharpness(full_enhanced).enhance(3.0)
    full_enhanced = full_enhanced.filter(ImageFilter.UnsharpMask(radius=2, percent=175, threshold=2))
    variants.append(full_enhanced)

    # 계기가 사진 중앙에 있는 일반적인 경우를 위한 넓은 문맥 확대본.
    # 자동 후보가 숫자 일부만 잡아도 온도·습도 화면 전체를 함께 확인할 수 있다.
    w, h = img.size
    if w >= 100 and h >= 100:
        context_crop = img.crop((int(w * 0.03), int(h * 0.02), int(w * 0.97), int(h * 0.86)))
        context_crop = ImageOps.autocontrast(context_crop, cutoff=1)
        context_crop = ImageEnhance.Contrast(context_crop).enhance(1.42)
        context_crop = ImageEnhance.Sharpness(context_crop).enhance(3.2)
        variants.append(context_crop)

    # 사진에서 LCD 가능성이 높은 영역을 자동 탐색해 크게 확대한다.
    lcd_crops = detect_heat_lcd_candidate_crops(img, max_candidates=2)
    for crop_index, crop in enumerate(lcd_crops):
        crop = ImageOps.autocontrast(crop, cutoff=1)
        crop = ImageEnhance.Contrast(crop).enhance(1.55)
        crop = ImageEnhance.Sharpness(crop).enhance(3.8)
        crop = crop.filter(ImageFilter.UnsharpMask(radius=2, percent=210, threshold=1))
        variants.append(crop)

        # 최우선 후보는 흑백 세그먼트 강화본도 추가한다.
        if crop_index == 0:
            gray = ImageOps.autocontrast(crop.convert("L"), cutoff=1)
            gray = ImageEnhance.Contrast(gray).enhance(1.85)
            gray = ImageEnhance.Sharpness(gray).enhance(3.4)
            variants.append(gray.convert("RGB"))

    data_urls = []
    seen_hashes = set()
    for variant in variants:
        try:
            thumb = variant.copy()
            thumb.thumbnail((160, 160))
            key_buf = io.BytesIO()
            thumb.save(key_buf, format="PNG")
            key = hashlib.sha256(key_buf.getvalue()).hexdigest()
            if key in seen_hashes:
                continue
            seen_hashes.add(key)

            data_urls.append(
                pil_image_to_data_url(
                    variant,
                    max_dim=3000,
                    min_dim=2200,
                    quality=95,
                )
            )
        except Exception:
            continue

    return data_urls[:6]


HEAT_EXTRACT_PROMPT = """이 이미지는 건설현장 체감온도 측정 기록 사진이다. 카카오톡 대화 캡쳐처럼
서로 다른 위치/시각의 측정 기록이 한 장에 여러 건 섞여 있을 수 있다. 각 기록을 구분해서 각각 하나의
JSON 객체로 만들어라.

각 기록에서 다음을 추출하라:
- 측정자: 다음 11명 중 이미지에서 확인되는 이름만 적는다: 김판식, 장경배, 박대우, 김종기, 송성태, 손만준, 조운제, 이용영, 방선혁, 공병대, 김명수. 직함은 제외한다. 명단 외 이름이거나 불확실하면 빈 문자열로 둔다.
- 측정위치: 사진 캡션이나 근처 텍스트에 명시된 현장/구역/동 이름
- 측정일자: 이미지에 날짜가 명시적으로 보이는 경우에만 YYYY-MM-DD 형식으로. 안 보이면 빈 문자열.
- 측정시간: 말풍선 옆 시각(오전/오후 표기 가능)을 24시간제 HH:MM으로 변환. 안 보이면 빈 문자열.
- 기온: 실제 온습도계 LCD의 위쪽 큰 숫자. 소수점과 ℃를 확인하고 숫자만 적는다.
- 습도: 실제 온습도계 LCD의 아래쪽 큰 숫자. % 표시와 가까운 숫자만 적는다.
- 체감온도: 계산기 화면 등에 명확히 숫자로 표시된 경우에만. 안 보이면 빈 문자열 (직접 계산하거나 추측하지 마라).
- 기온확신도: LCD 기온 판독 확신도를 0~1 숫자로 표시한다.
- 습도확신도: LCD 습도 판독 확신도를 0~1 숫자로 표시한다.

실제 측정기 판독 규칙:
- 사진 설명글이나 채팅 문구에 적힌 숫자보다 온습도계 LCD에 실제 표시된 값을 우선한다.
- 일반적인 온습도계는 위쪽 큰 숫자가 기온(℃), 아래쪽 큰 숫자가 습도(%)다.
- 화면 왼쪽 아래의 작은 3~4자리 숫자는 시각이나 보조표시일 수 있으므로 습도로 읽지 마라.
- 예를 들어 위쪽에 30.0℃, 아래쪽에 63%가 보이면 기온 30.0, 습도 63으로 읽는다.
- 기온 소수점이 흐리더라도 300으로 쓰지 말고 LCD 배치와 ℃ 표시를 확인해 30.0처럼 판독한다.
- 습도 숫자 뒤에 % 표시가 붙어 있는지 확인한다.

판독 정확도 관련 유의사항:
- LCD 숫자는 세그먼트 특성상 2↔7, 3↔8, 1↔7, 6↔8, 0↔8이 헷갈릴 수 있다.
  각 자릿수의 켜진 막대를 다시 비교해 판독하라.
- 반사광, 원거리 촬영, 저화질 때문에 값이 흐려도 먼저 자동 탐색·확대된 LCD 후보 영역을 원본과 대조해 실제 표시값을 읽어보라.
- 한국 건설현장의 여름철 기온은 대체로 15~40℃, 습도는 20~100% 범위다. 범위를 크게 벗어나면
  오독 가능성을 다시 확인하되, 화면에 분명히 표시된 값이라면 읽은 그대로 적어라.
- 끝까지 구분할 수 없는 값만 빈 문자열("")로 남겨라. 단순히 조금 흐리다는 이유로 바로 포기하지 마라.

아래 JSON 배열 형식으로만 출력하라. 설명, 코드블록, 다른 텍스트를 절대 포함하지 마라.
[
  {"측정자":"", "측정위치":"", "측정일자":"", "측정시간":"", "기온":"", "습도":"", "체감온도":"", "기온확신도":0.0, "습도확신도":0.0}
]

기록이 하나도 없으면 빈 배열 []만 출력하라."""


HEAT_METER_RETRY_PROMPT = """여러 입력 이미지는 모두 동일한 원본 사진 또는 프로그램이 자동 탐색해 크게 확대한 LCD 후보 영역이다.
서로 다른 사진으로 세지 말고, 같은 온습도계를 중복해서 기록하지 마라.

사진 속 실제 디지털 온습도계 LCD를 직접 판독하라. 채팅 문구나 주변 텍스트가 아니라 LCD 실측값이 기준이다.
- 위쪽 큰 숫자 + ℃ 표시 = 기온
- 아래쪽 큰 숫자 + % 표시 = 습도
- 왼쪽 아래의 작은 시각/보조 숫자는 습도가 아니다.
- 소수점을 반드시 확인한다. 예: 30.0℃를 300 또는 30으로 잘못 읽지 마라.
- 7세그먼트 숫자의 2/7, 3/8, 1/7, 6/8, 0/8을 켜진 막대 모양으로 비교한다.
- 원본, 색상 강화본, 흑백 강화본, 확대본을 서로 대조해 가장 일치하는 값을 선택한다.
- 한 원본 안에 측정기가 여러 개면 화면 위에서 아래, 왼쪽에서 오른쪽 순서로 각각 반환한다.
- 끝까지 판독 불가능한 항목만 빈 문자열로 둔다. 평균값이나 일반적인 값을 추측해서 넣지 마라.

다음 JSON 배열만 출력하라.
[
  {"기온":"", "습도":"", "기온확신도":0.0, "습도확신도":0.0}
]

측정기를 찾지 못하면 []만 출력하라."""


def _extract_openai_output_text(data: dict) -> str:
    text = ""
    if data.get("output_text"):
        return str(data.get("output_text") or "")
    for item in data.get("output", []):
        for content in item.get("content", []):
            if content.get("type") == "output_text":
                text += str(content.get("text", "") or "")
    return text


def _safe_confidence(value, default: float = 0.0) -> float:
    try:
        confidence = float(value)
    except (TypeError, ValueError):
        return default
    return max(0.0, min(1.0, confidence))


def parse_heat_meter_retry_response(text: str) -> List[dict]:
    cleaned = str(text or "").replace("```json", "").replace("```", "").strip()
    try:
        parsed = json.loads(cleaned)
    except Exception:
        cleaned = re.sub(r"[\x00-\x1F]+", " ", cleaned)
        match = re.search(r"\[.*\]", cleaned, re.S)
        parsed = json.loads(match.group(0)) if match else []

    if not isinstance(parsed, list):
        parsed = [parsed]

    results = []
    for item in parsed:
        if not isinstance(item, dict):
            continue
        results.append({
            "기온": str(item.get("기온", "") or "").strip(),
            "습도": str(item.get("습도", "") or "").strip(),
            "기온확신도": _safe_confidence(item.get("기온확신도"), 0.0),
            "습도확신도": _safe_confidence(item.get("습도확신도"), 0.0),
        })
    return results


def extract_heat_meter_values_with_gpt(api_key: str, image_path: str) -> List[dict]:
    """흐리거나 작은 실제 측정기 LCD를 확대·보정 이미지로 한 번 더 판독."""
    image_urls = make_heat_meter_vision_images(image_path)
    if not image_urls:
        return []

    content = [{"type": "input_text", "text": HEAT_METER_RETRY_PROMPT}]
    content.extend({"type": "input_image", "image_url": url} for url in image_urls)

    headers = {
        "Authorization": f"Bearer {api_key.strip()}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": "gpt-4o-mini",
        "input": [{"role": "user", "content": content}],
    }

    response = requests.post(
        "https://api.openai.com/v1/responses",
        headers=headers,
        json=payload,
        timeout=75,
    )
    if response.status_code != 200:
        return []

    return parse_heat_meter_retry_response(_extract_openai_output_text(response.json()))


def heat_record_needs_meter_retry(record: dict) -> bool:
    temp = parse_heat_number(record.get("기온"), "기온")
    humidity = parse_heat_number(record.get("습도"), "습도")
    temp_conf = _safe_confidence(record.get("_기온확신도"), 0.0)
    hum_conf = _safe_confidence(record.get("_습도확신도"), 0.0)

    if temp is None or humidity is None:
        return True
    if temp_conf < 0.82 or hum_conf < 0.82:
        return True
    return False


def _replace_meter_field(record: dict, retry: dict, field: str, direct_photo: bool = False):
    confidence_key = f"_{field}확신도"
    retry_confidence_key = f"{field}확신도"

    old_value = parse_heat_number(record.get(field), field)
    new_value = parse_heat_number(retry.get(field), field)
    if new_value is None:
        return

    old_conf = _safe_confidence(record.get(confidence_key), 0.0)
    new_conf = _safe_confidence(retry.get(retry_confidence_key), 0.0)

    should_replace = False
    if old_value is None:
        should_replace = True
    elif old_conf < 0.82 and new_conf >= old_conf:
        should_replace = True
    elif direct_photo and new_conf >= 0.90 and new_conf > old_conf:
        # 측정기 단독 사진은 재판독 결과가 더 확실하면 기존 값도 교정한다.
        should_replace = True

    if should_replace:
        record[field] = str(int(new_value)) if field == "습도" and new_value.is_integer() else str(new_value)
        record[confidence_key] = new_conf
        record["_meter_retry_used"] = True


def merge_heat_meter_retry(records: List[dict], retry_values: List[dict]) -> List[dict]:
    """1차 구조화 결과에 LCD 전용 재판독값을 안전하게 합침."""
    if not retry_values:
        return records

    # 온습도계 사진만 단독으로 올린 경우에도 기록 1건을 생성한다.
    if not records:
        return [{
            "측정자": "",
            "측정위치": "",
            "측정일자": "",
            "측정시간": "",
            "기온": item.get("기온", ""),
            "습도": item.get("습도", ""),
            "체감온도": "",
            "_기온확신도": _safe_confidence(item.get("기온확신도"), 0.0),
            "_습도확신도": _safe_confidence(item.get("습도확신도"), 0.0),
            "_meter_retry_used": True,
        } for item in retry_values]

    direct_photo = (
        len(records) == 1
        and not records[0].get("측정자")
        and not records[0].get("측정위치")
        and not records[0].get("측정일자")
        and not records[0].get("측정시간")
    )

    if len(records) == len(retry_values):
        pairs = list(zip(range(len(records)), retry_values))
    else:
        retry_indexes = [
            idx for idx, record in enumerate(records)
            if heat_record_needs_meter_retry(record)
        ]
        pairs = list(zip(retry_indexes, retry_values))

    for record_index, retry in pairs:
        record = records[record_index]
        _replace_meter_field(record, retry, "기온", direct_photo=direct_photo)
        _replace_meter_field(record, retry, "습도", direct_photo=direct_photo)

        # 기온 또는 습도가 재판독으로 바뀌었다면 기존 체감온도는 공식으로 다시 산출하도록 비운다.
        if record.get("_meter_retry_used"):
            record["체감온도"] = ""

    return records


def extract_heat_records_with_gpt(api_key: str, image_path: str) -> List[dict]:
    """사진 1장에서 기록을 추출하고, 흐린 LCD는 확대·보정 후 자동 재판독."""
    url = "https://api.openai.com/v1/responses"
    data_url = image_to_data_url(image_path)

    headers = {
        "Authorization": f"Bearer {api_key.strip()}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": "gpt-4o-mini",
        "input": [
            {
                "role": "user",
                "content": [
                    {"type": "input_text", "text": HEAT_EXTRACT_PROMPT},
                    {"type": "input_image", "image_url": data_url},
                ],
            }
        ],
    }

    response = requests.post(url, headers=headers, json=payload, timeout=60)
    if response.status_code != 200:
        raise Exception(f"API Error: {response.text}")

    records = parse_heat_gpt_response(_extract_openai_output_text(response.json()))

    # 값이 비거나 확신도가 낮은 경우에만 LCD 전용 2차 판독을 실행한다.
    # 온습도계 단독 사진도 실제 LCD 값을 재확인한다.
    direct_photo = (
        len(records) == 1
        and not records[0].get("측정자")
        and not records[0].get("측정위치")
        and not records[0].get("측정일자")
        and not records[0].get("측정시간")
    )
    needs_retry = not records or any(heat_record_needs_meter_retry(r) for r in records) or direct_photo

    if needs_retry:
        retry_values = extract_heat_meter_values_with_gpt(api_key, image_path)
        records = merge_heat_meter_retry(records, retry_values)

    return records


def parse_heat_gpt_response(text: str) -> List[dict]:
    """GPT 응답 텍스트를 JSON 배열로 파싱 (코드블록/잡텍스트 방어)."""
    cleaned = str(text or "").replace("```json", "").replace("```", "").strip()
    try:
        parsed = json.loads(cleaned)
    except Exception:
        cleaned = re.sub(r"[\x00-\x1F]+", " ", cleaned)
        match = re.search(r"\[.*\]", cleaned, re.S)
        parsed = json.loads(match.group(0)) if match else []

    if not isinstance(parsed, list):
        parsed = [parsed]

    records = []
    for item in parsed:
        if not isinstance(item, dict):
            continue
        record = {
            "측정자": str(item.get("측정자", "") or "").strip(),
            "측정위치": str(item.get("측정위치", "") or "").strip(),
            "측정일자": str(item.get("측정일자", "") or "").strip(),
            "측정시간": str(item.get("측정시간", "") or "").strip(),
            "기온": str(item.get("기온", "") or "").strip(),
            "습도": str(item.get("습도", "") or "").strip(),
            "체감온도": str(item.get("체감온도", "") or "").strip(),
            "_기온확신도": _safe_confidence(item.get("기온확신도"), 0.0),
            "_습도확신도": _safe_confidence(item.get("습도확신도"), 0.0),
        }
        records.append(record)
    return records


def load_heat_template_sheet():
    """템플릿 파일을 열어 원본 시트를 반환. 없으면 None."""
    if not os.path.exists(HEAT_TEMPLATE_XLSX):
        return None
    try:
        wb = openpyxl.load_workbook(HEAT_TEMPLATE_XLSX)
        return wb[wb.sheetnames[0]]
    except Exception:
        return None


def copy_template_sheet(template_ws, target_wb, new_title: str):
    """템플릿 시트를 다른 워크북으로 값+서식+병합+열너비까지 복제."""
    new_ws = target_wb.create_sheet(title=new_title[:31])

    for row in template_ws.iter_rows():
        for cell in row:
            new_cell = new_ws.cell(row=cell.row, column=cell.column, value=cell.value)
            if cell.has_style:
                new_cell.font = copy(cell.font)
                new_cell.border = copy(cell.border)
                new_cell.fill = copy(cell.fill)
                new_cell.number_format = cell.number_format
                new_cell.protection = copy(cell.protection)
                new_cell.alignment = copy(cell.alignment)

    for mc in template_ws.merged_cells.ranges:
        new_ws.merge_cells(str(mc))

    for col_letter, dim in template_ws.column_dimensions.items():
        new_ws.column_dimensions[col_letter].width = dim.width

    for row_idx, dim in template_ws.row_dimensions.items():
        new_ws.row_dimensions[row_idx].height = dim.height

    return new_ws


def date_to_mmdd(date_str: str) -> str:
    d = re.sub(r"[./]", "-", str(date_str or "").strip())
    m = re.match(r"(\d{4})-(\d{1,2})-(\d{1,2})", d)
    if m:
        return f"{int(m.group(2)):02d}{int(m.group(3)):02d}"
    digits = re.sub(r"\D", "", d)
    return digits[-4:] if len(digits) >= 4 else "0000"


def get_or_create_heat_sheet(wb, template_ws, location: str, date_str: str):
    """같은 (장소, 날짜)면 기존 시트 재사용, 표기 차이는 유사도로 병합. 없으면 템플릿을 복제해 새로 생성."""
    mmdd = date_to_mmdd(date_str)

    candidates = []
    for name in wb.sheetnames:
        m = re.match(r"^(.*)_(\d{4})(?:_\d+)?$", name)
        if m and m.group(2) == mmdd:
            candidates.append((name, m.group(1)))

    loc_norm = normalize_for_match(location)
    best_name, best_score = None, 0.0
    for name, prefix in candidates:
        score = SequenceMatcher(None, normalize_for_match(prefix), loc_norm).ratio()
        if score > best_score:
            best_score = score
            best_name = name

    if best_score >= 0.75:
        return wb[best_name]

    safe_loc = re.sub(r'[\\/*?:\[\]]', "_", location).strip()[:20] or "미지정"
    base_name = f"{safe_loc}_{mmdd}"[:31]
    new_name = base_name
    i = 2
    while new_name in wb.sheetnames:
        new_name = f"{base_name}_{i}"[:31]
        i += 1

    ws = copy_template_sheet(template_ws, wb, new_name)

    orig_b3 = template_ws["B3"].value or ""
    new_b3 = re.sub(r"(측정일자\s*:)\s*", rf"\1 {date_str}   ", orig_b3, count=1)
    new_b3 = re.sub(r"(측정장소\s*:)\s*", rf"\1 {location}   ", new_b3, count=1)
    ws["B3"] = new_b3

    return ws


def is_empty_heat_slot(ws, template_ws, row: int) -> bool:
    """템플릿 원본과 C~H 값이 같거나 모두 비어 있으면 아직 사용하지 않은 행으로 판단."""
    current_values = [ws.cell(row=row, column=col).value for col in range(3, 9)]
    template_values = [template_ws.cell(row=row, column=col).value for col in range(3, 9)]

    if all(v in (None, "") for v in current_values):
        return True

    return all(
        current in (None, "") if template in (None, "") else current == template
        for current, template in zip(current_values, template_values)
    )


def find_empty_heat_slot(ws, template_ws) -> Optional[int]:
    """NO 1~9(6~14행) 중 아직 안 쓴 첫 번째 칸을 찾음. 다 찼으면 None.

    기존에는 기온(D열)이 비어 있으면 이미 기록된 행도 빈 행으로 오인할 수 있었다.
    이제 시간·습도·체감온도·측정자·비고까지 포함한 C~H 전체를 확인한다.
    """
    for row in range(6, 15):
        if is_empty_heat_slot(ws, template_ws, row):
            return row
    return None


def heat_time_to_minutes(t) -> Optional[int]:
    s = str(t or "").strip()
    m = re.match(r"^(\d{1,2}):(\d{2})$", s)
    if not m:
        m = re.match(r"^(\d{1,2})\s*시\s*(\d{1,2})\s*분$", s)
    if not m:
        return None
    return int(m.group(1)) * 60 + int(m.group(2))


def format_heat_time_display(t: str) -> str:
    """'08:55' -> '08시 55분'. 이미 'OO시 OO분' 형식이면 그대로 둠. 파싱 안 되면 원문 유지."""
    s = str(t or "").strip()
    m = re.match(r"^(\d{1,2}):(\d{2})$", s)
    if m:
        return f"{int(m.group(1)):02d}시 {int(m.group(2)):02d}분"
    if re.match(r"^\d{1,2}\s*시\s*\d{1,2}\s*분$", s):
        return s
    return s


HEAT_ALLOWED_MEASURERS = [
    "김판식", "장경배", "박대우", "김종기", "송성태", "손만준",
    "조운제", "이용영", "방선혁", "공병대", "김명수",
]
HEAT_MEASURER_SUFFIXES = ["대원", "반장", "조장", "소장", "관리자", "주임", "팀장", "과장", "부장", "님"]


def match_allowed_measurer(name: str) -> str:
    """OCR 이름을 지정된 11명 중 한 명으로만 보수적으로 매칭. 불확실하면 빈 문자열."""
    raw = str(name or "").strip()
    if not raw:
        return ""

    compact = re.sub(r"[^가-힣]", "", raw)
    for suffix in HEAT_MEASURER_SUFFIXES:
        suffix_compact = re.sub(r"[^가-힣]", "", suffix)
        if compact.endswith(suffix_compact):
            compact = compact[:-len(suffix_compact)]
            break

    for allowed in HEAT_ALLOWED_MEASURERS:
        if allowed in compact or compact in allowed and len(compact) >= 2:
            return allowed

    best_name, best_score, second_score = "", 0.0, 0.0
    for allowed in HEAT_ALLOWED_MEASURERS:
        score = SequenceMatcher(None, compact, allowed).ratio()
        if score > best_score:
            second_score = best_score
            best_score = score
            best_name = allowed
        elif score > second_score:
            second_score = score

    # 다른 이름을 억지로 특정 대원으로 바꾸지 않도록 높은 기준과 점수 차이를 함께 요구한다.
    if best_score >= 0.74 and best_score - second_score >= 0.08:
        return best_name
    return ""


def format_measurer(name: str) -> str:
    """측정자 표기는 지정된 11명만 허용하며 항상 '이름 대원'으로 통일."""
    matched = match_allowed_measurer(name)
    return f"{matched} 대원" if matched else ""


def is_implausible_value(temp, humidity) -> bool:
    """OCR/판독 오류일 가능성이 높은 비현실적 수치인지 확인 (자동 대체는 하지 않고 표시만 함)."""
    try:
        if temp not in (None, "") and not (10 <= float(temp) <= 40):
            return True
    except (TypeError, ValueError):
        pass
    try:
        if humidity not in (None, "") and not (0 <= float(humidity) <= 100):
            return True
    except (TypeError, ValueError):
        pass
    return False


HEAT_EXPECTED_SLOTS = [("09시대", 8 * 60, 10 * 60), ("11시대", 10 * 60, 12 * 60),
                        ("13시대", 12 * 60, 14 * 60), ("15시대", 14 * 60, 16 * 60)]


def get_slot_coverage_text(ws) -> str:
    """하루 기본 4회(9/11/13/15시 전후) 측정 슬롯 중 실제 기록이 있는 슬롯을 체크 표시로 보여줌."""
    recorded_minutes = []
    for r in range(6, 15):
        mins = heat_time_to_minutes(ws.cell(row=r, column=3).value)
        if mins is not None:
            recorded_minutes.append(mins)

    parts = []
    for label, start, end in HEAT_EXPECTED_SLOTS:
        covered = any(start <= m < end for m in recorded_minutes)
        parts.append(f"{label} {'✅' if covered else '❌'}")
    return "  ".join(parts)


def _to_number(value):
    try:
        return float(value)
    except (TypeError, ValueError):
        return value


HEAT_NUMERIC_FIELDS = {
    "기온": 4,
    "습도": 5,
    "체감온도": 6,
}

HEAT_AUTO_NOTE_PREFIXES = (
    "평균치 자동기입",
    "체감온도 자동산출",
    "평균 산출 불가",
)


def parse_heat_number(value, field_name: str = "") -> Optional[float]:
    """'31.2℃', '65%' 같은 값도 숫자로 정리하고 평균 산출에 쓸 수 있는 범위인지 확인."""
    if value in (None, ""):
        return None

    s = str(value).strip().replace(",", "")
    m = re.search(r"-?\d+(?:\.\d+)?", s)
    if not m:
        return None

    try:
        number = float(m.group(0))
    except (TypeError, ValueError):
        return None

    if field_name == "기온" and not (10 <= number <= 40):
        return None
    if field_name == "습도" and not (0 <= number <= 100):
        return None
    if field_name == "체감온도" and not (10 <= number <= 60):
        return None
    return number


def round_heat_average(value: float) -> float:
    return round(float(value), 1)


def get_heat_field_averages_from_sheet(ws, exclude_row: Optional[int] = None) -> dict:
    """같은 장소·날짜 시트의 정상 기록 평균을 구함. 유효한 측정시간이 있는 행만 사용."""
    values = {field: [] for field in HEAT_NUMERIC_FIELDS}

    for row_idx in range(6, 15):
        if exclude_row is not None and row_idx == exclude_row:
            continue
        if heat_time_to_minutes(ws.cell(row=row_idx, column=3).value) is None:
            continue

        for field, col in HEAT_NUMERIC_FIELDS.items():
            number = parse_heat_number(ws.cell(row=row_idx, column=col).value, field)
            if number is not None:
                values[field].append(number)

    return {
        field: round_heat_average(sum(nums) / len(nums))
        for field, nums in values.items()
        if nums
    }


def get_heat_field_averages_from_workbook_date(wb, date_str: str, exclude_sheet: str = "") -> dict:
    """같은 장소 시트에 평균 자료가 없을 때, 같은 날짜의 다른 장소 기록 평균을 보조로 사용."""
    mmdd = date_to_mmdd(date_str)
    values = {field: [] for field in HEAT_NUMERIC_FIELDS}

    for sheet_name in wb.sheetnames:
        if sheet_name == exclude_sheet:
            continue
        if not re.search(rf"_{re.escape(mmdd)}(?:_\d+)?$", sheet_name):
            continue

        ws = wb[sheet_name]
        for row_idx in range(6, 15):
            if heat_time_to_minutes(ws.cell(row=row_idx, column=3).value) is None:
                continue
            for field, col in HEAT_NUMERIC_FIELDS.items():
                number = parse_heat_number(ws.cell(row=row_idx, column=col).value, field)
                if number is not None:
                    values[field].append(number)

    return {
        field: round_heat_average(sum(nums) / len(nums))
        for field, nums in values.items()
        if nums
    }


def fill_missing_heat_values(
    row: dict,
    primary_averages: dict,
    fallback_averages: Optional[dict] = None,
) -> Tuple[List[str], bool, List[str]]:
    """빈 기온·습도·체감온도를 평균으로 보완.

    우선순위:
    1) 같은 장소·같은 날짜 평균
    2) 같은 날짜 다른 장소 평균
    3) 기온·습도가 확보되면 체감온도 공식 계산
    """
    fallback_averages = fallback_averages or {}
    averaged_fields = list(row.get("_average_filled", []))
    unavailable_fields = []
    auto_calculated = bool(row.get("_heat_index_calculated", False))

    for field in ("기온", "습도"):
        if parse_heat_number(row.get(field), field) is not None:
            continue
        average_value = primary_averages.get(field, fallback_averages.get(field))
        if average_value is not None:
            row[field] = str(round_heat_average(average_value))
            if field not in averaged_fields:
                averaged_fields.append(field)
        else:
            unavailable_fields.append(field)

    if parse_heat_number(row.get("체감온도"), "체감온도") is None:
        temp = parse_heat_number(row.get("기온"), "기온")
        humidity = parse_heat_number(row.get("습도"), "습도")
        if temp is not None and humidity is not None:
            row["체감온도"] = str(calc_heat_index(temp, humidity))
            auto_calculated = True
        else:
            average_value = primary_averages.get("체감온도", fallback_averages.get("체감온도"))
            if average_value is not None:
                row["체감온도"] = str(round_heat_average(average_value))
                if "체감온도" not in averaged_fields:
                    averaged_fields.append("체감온도")
            else:
                unavailable_fields.append("체감온도")

    row["_average_filled"] = averaged_fields
    row["_heat_index_calculated"] = auto_calculated
    return averaged_fields, auto_calculated, unavailable_fields


def prepare_heat_records_for_average(records: List[dict]) -> List[dict]:
    """한 사진에서 추출된 기록끼리 먼저 평균 보완해, 첫 저장 건도 뒤쪽 정상값을 활용할 수 있게 함."""
    today = datetime.now().strftime("%Y-%m-%d")
    now_time = datetime.now().strftime("%H:%M")

    for rec in records:
        raw_person = str(rec.get("측정자", "") or "").strip()
        formatted_person = format_measurer(raw_person)
        rec["측정자"] = formatted_person
        rec["_measurer_unrecognized"] = bool(raw_person and not formatted_person)
        if not rec.get("측정위치"):
            rec["측정위치"] = "미지정"
        if not rec.get("측정일자"):
            rec["측정일자"] = today
        if not rec.get("측정시간"):
            rec["측정시간"] = now_time

    grouped = {}
    date_grouped = {}
    for rec in records:
        key = (normalize_for_match(rec.get("측정위치", "")), rec.get("측정일자", ""))
        grouped.setdefault(key, []).append(rec)
        date_grouped.setdefault(rec.get("측정일자", ""), []).append(rec)

    def averages_from_records(group: List[dict]) -> dict:
        result = {}
        for field in HEAT_NUMERIC_FIELDS:
            nums = [parse_heat_number(item.get(field), field) for item in group]
            nums = [n for n in nums if n is not None]
            if nums:
                result[field] = round_heat_average(sum(nums) / len(nums))
        return result

    group_averages = {key: averages_from_records(group) for key, group in grouped.items()}
    date_averages = {key: averages_from_records(group) for key, group in date_grouped.items()}

    for rec in records:
        key = (normalize_for_match(rec.get("측정위치", "")), rec.get("측정일자", ""))
        fill_missing_heat_values(
            rec,
            group_averages.get(key, {}),
            date_averages.get(rec.get("측정일자", ""), {}),
        )

    return records


def clean_heat_auto_notes(existing_note: str) -> List[str]:
    """수정 반영 시 이전 자동 평균/산출 메모만 제거하고 간격초과 등 다른 메모는 유지."""
    parts = [part.strip() for part in str(existing_note or "").split(" / ") if part.strip()]
    return [
        part for part in parts
        if not any(part.startswith(prefix) for prefix in HEAT_AUTO_NOTE_PREFIXES)
    ]


def build_heat_auto_notes(row: dict, unavailable_fields: Optional[List[str]] = None) -> List[str]:
    notes = []
    averaged_fields = list(dict.fromkeys(row.get("_average_filled", [])))
    if averaged_fields:
        details = []
        for field in averaged_fields:
            value = row.get(field, "")
            unit = "%" if field == "습도" else "℃"
            details.append(f"{field} {value}{unit}")
        notes.append("평균치 자동기입(" + ", ".join(details) + ")")

    if row.get("_heat_index_calculated"):
        notes.append("체감온도 자동산출(기온·습도 기준)")

    if unavailable_fields:
        unique_fields = list(dict.fromkeys(unavailable_fields))
        notes.append("평균 산출 불가(" + ", ".join(unique_fields) + ")")
    return notes



def ensure_heat_history_storage():
    os.makedirs(HEAT_EXPORT_HISTORY_DIR, exist_ok=True)


def _empty_heat_upload_history() -> dict:
    return {
        "processed_files": {},
        "batches": [],
    }


def load_heat_upload_history() -> dict:
    """파일 해시와 시간대별 완료본 목록을 영구 저장한 JSON을 읽음."""
    ensure_heat_history_storage()
    if not os.path.exists(HEAT_UPLOAD_HISTORY_FILE):
        return _empty_heat_upload_history()

    try:
        with open(HEAT_UPLOAD_HISTORY_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, dict):
            return _empty_heat_upload_history()
        data.setdefault("processed_files", {})
        data.setdefault("batches", [])
        return data
    except Exception:
        return _empty_heat_upload_history()


def save_heat_upload_history(history: dict):
    """기록 손상을 줄이기 위해 임시 파일에 쓴 뒤 원자적으로 교체."""
    ensure_heat_history_storage()
    temp_path = HEAT_UPLOAD_HISTORY_FILE + ".tmp"
    with open(temp_path, "w", encoding="utf-8") as f:
        json.dump(history, f, ensure_ascii=False, indent=2)
    os.replace(temp_path, HEAT_UPLOAD_HISTORY_FILE)


def uploaded_file_hash(uploaded_file) -> Tuple[str, bytes]:
    """파일명과 무관하게 실제 바이트 기준 SHA-256 해시를 계산."""
    file_bytes = uploaded_file.getvalue()
    return hashlib.sha256(file_bytes).hexdigest(), file_bytes


def _parse_history_datetime(value: str) -> Optional[datetime]:
    try:
        return datetime.strptime(str(value or ""), "%Y-%m-%d %H:%M:%S")
    except Exception:
        return None


def reserve_heat_upload(file_hash: str, original_name: str, file_size: int) -> Tuple[bool, dict]:
    """새 파일이면 처리 예약. 이미 처리된 파일이면 False와 기존 메타데이터 반환."""
    now = datetime.now()
    with _HEAT_HISTORY_LOCK:
        history = load_heat_upload_history()
        existing = history["processed_files"].get(file_hash)

        if existing:
            status = str(existing.get("status", "completed"))
            started_at = _parse_history_datetime(existing.get("started_at", ""))
            stale_processing = (
                status == "processing"
                and started_at is not None
                and (now - started_at).total_seconds() > HEAT_PROCESSING_STALE_SECONDS
            )

            # 오류 건은 다음 세션에서 재시도 가능, 오래 멈춘 processing도 다시 처리.
            if status != "error" and not stale_processing:
                return False, existing

        reserved = {
            "file_hash": file_hash,
            "original_name": original_name,
            "file_size": int(file_size or 0),
            "status": "processing",
            "started_at": now.strftime("%Y-%m-%d %H:%M:%S"),
            "processed_at": "",
            "records_count": 0,
            "new_records_count": 0,
            "duplicate_records_count": 0,
            "sheet_names": [],
            "batch_id": "",
            "error": "",
        }
        history["processed_files"][file_hash] = reserved
        save_heat_upload_history(history)
        return True, reserved


def finalize_heat_upload(
    file_hash: str,
    status: str,
    entries: List[dict],
    error: str = "",
):
    successful_entries = [e for e in entries if not e.get("error")]
    new_entries = [e for e in successful_entries if not e.get("duplicate")]
    duplicate_entries = [e for e in successful_entries if e.get("duplicate")]
    sheet_names = sorted({e.get("sheet", "") for e in successful_entries if e.get("sheet")})

    with _HEAT_HISTORY_LOCK:
        history = load_heat_upload_history()
        meta = history["processed_files"].setdefault(file_hash, {"file_hash": file_hash})
        meta.update({
            "status": status,
            "processed_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "records_count": len(successful_entries),
            "new_records_count": len(new_entries),
            "duplicate_records_count": len(duplicate_entries),
            "sheet_names": sheet_names,
            "error": str(error or ""),
        })
        save_heat_upload_history(history)


def create_heat_batch_snapshot(new_file_infos: List[dict], new_entries: List[dict]) -> Optional[dict]:
    """이번 업로드에서 정상 인식된 기록이 있으면 누적 엑셀 스냅샷 1개 생성.

    완료 파일을 삭제한 뒤 같은 사진을 다시 올린 경우에는 기존 행과 중복이더라도
    전체 누적 대장을 다시 내려받을 수 있도록 완료 파일을 재생성한다.
    """
    successful_entries = [e for e in new_entries if not e.get("error")]
    actual_new_entries = [e for e in successful_entries if not e.get("duplicate")]
    reused_entries = [e for e in successful_entries if e.get("duplicate")]
    if not new_file_infos or not successful_entries or not os.path.exists(HEAT_LOG_FILE):
        return None

    file_hashes = sorted({info["file_hash"] for info in new_file_infos})
    batch_signature = hashlib.sha256("|".join(file_hashes).encode("utf-8")).hexdigest()

    with _HEAT_HISTORY_LOCK:
        history = load_heat_upload_history()
        for batch in history.get("batches", []):
            if batch.get("batch_signature") == batch_signature:
                return batch

        now = datetime.now()
        stamp = now.strftime("%Y%m%d_%H%M%S")
        batch_id = f"{stamp}_{batch_signature[:8]}"
        filename = f"체감온도측정_누적대장_{stamp}.xlsx"
        snapshot_path = os.path.join(HEAT_EXPORT_HISTORY_DIR, filename)

        suffix = 2
        while os.path.exists(snapshot_path):
            filename = f"체감온도측정_누적대장_{stamp}_{suffix}.xlsx"
            snapshot_path = os.path.join(HEAT_EXPORT_HISTORY_DIR, filename)
            suffix += 1

        shutil.copy2(HEAT_LOG_FILE, snapshot_path)

        batch = {
            "batch_id": batch_id,
            "batch_signature": batch_signature,
            "created_at": now.strftime("%Y-%m-%d %H:%M:%S"),
            "snapshot_filename": filename,
            "source_files": [info.get("original_name", "") for info in new_file_infos],
            "source_file_count": len(new_file_infos),
            "new_records_count": len(actual_new_entries),
            "reused_records_count": len(reused_entries),
            "recognized_records_count": len(successful_entries),
            "sheet_names": sorted({e.get("sheet", "") for e in successful_entries if e.get("sheet")}),
        }
        history.setdefault("batches", []).append(batch)

        for info in new_file_infos:
            meta = history["processed_files"].get(info["file_hash"])
            if meta is not None:
                meta["batch_id"] = batch_id

        save_heat_upload_history(history)
        return batch


def get_heat_completed_batches() -> List[dict]:
    history = load_heat_upload_history()
    batches = []
    for batch in history.get("batches", []):
        filename = batch.get("snapshot_filename", "")
        path = os.path.join(HEAT_EXPORT_HISTORY_DIR, filename) if filename else ""
        if path and os.path.exists(path):
            item = dict(batch)
            item["snapshot_path"] = path
            batches.append(item)
    return sorted(batches, key=lambda x: x.get("created_at", ""), reverse=True)


def delete_heat_completed_batch(batch_id: str) -> Tuple[bool, str]:
    """시간대별 완료 파일과 목록을 삭제하고 해당 원본 사진의 재생성을 허용.

    전체 누적 대장 행은 유지한다. 다만 이 완료 파일에 연결된 파일 해시만 제거해
    같은 사진을 다시 올렸을 때 OCR 확인 후 완료 엑셀을 다시 만들 수 있게 한다.
    """
    batch_id = str(batch_id or "").strip()
    if not batch_id:
        return False, "삭제할 완료 파일 식별값이 없습니다."

    with _HEAT_HISTORY_LOCK:
        history = load_heat_upload_history()
        batches = history.get("batches", [])
        target = next((b for b in batches if str(b.get("batch_id", "")) == batch_id), None)

        if target is None:
            return False, "이미 삭제되었거나 완료 파일을 찾을 수 없습니다."

        filename = os.path.basename(str(target.get("snapshot_filename", "") or ""))
        if filename:
            snapshot_path = os.path.join(HEAT_EXPORT_HISTORY_DIR, filename)
            try:
                if os.path.isfile(snapshot_path):
                    os.remove(snapshot_path)
            except Exception as e:
                return False, f"완료 파일 삭제 실패: {e}"

        history["batches"] = [
            b for b in batches
            if str(b.get("batch_id", "")) != batch_id
        ]

        # 삭제된 완료본에 포함된 사진은 다시 업로드해 완료 파일을 재생성할 수 있도록
        # 파일 단위 처리 이력만 제거한다. 전체 누적 대장의 측정 행은 삭제하지 않는다.
        processed_files = history.get("processed_files", {})
        removable_hashes = [
            file_hash for file_hash, meta in processed_files.items()
            if str(meta.get("batch_id", "")) == batch_id
        ]
        for file_hash in removable_hashes:
            processed_files.pop(file_hash, None)

        save_heat_upload_history(history)

    return True, "완료 파일을 삭제했습니다. 같은 사진을 다시 올리면 완료 엑셀을 재생성할 수 있습니다."


def _same_heat_number(a, b, field_name: str) -> bool:
    na = parse_heat_number(a, field_name)
    nb = parse_heat_number(b, field_name)
    if na is None or nb is None:
        return False
    return abs(na - nb) < 0.05


def find_duplicate_heat_row(ws, row: dict) -> Optional[int]:
    """파일 바이트가 달라도 장소·날짜·시간과 측정값이 같은 기록이면 중복으로 판단."""
    target_minutes = heat_time_to_minutes(row.get("측정시간"))
    if target_minutes is None:
        return None

    target_person = normalize_for_match(format_measurer(row.get("측정자", "")))

    for row_idx in range(6, 15):
        existing_minutes = heat_time_to_minutes(ws.cell(row=row_idx, column=3).value)
        if existing_minutes != target_minutes:
            continue

        score = 0
        if _same_heat_number(ws.cell(row=row_idx, column=4).value, row.get("기온"), "기온"):
            score += 1
        if _same_heat_number(ws.cell(row=row_idx, column=5).value, row.get("습도"), "습도"):
            score += 1
        if _same_heat_number(ws.cell(row=row_idx, column=6).value, row.get("체감온도"), "체감온도"):
            score += 1

        existing_person = normalize_for_match(ws.cell(row=row_idx, column=7).value)
        if target_person and existing_person and target_person == existing_person:
            score += 1

        # 동일 시간에 2개 이상의 핵심 값이 일치하면 같은 측정 기록으로 봄.
        if score >= 2:
            return row_idx

    return None


def render_heat_completed_file_list():
    """시간대별 완료 엑셀을 다운로드하거나 목록에서 개별 삭제."""
    batches = get_heat_completed_batches()
    st.markdown("##### 시간대별 완료 파일")

    if not batches:
        st.caption("아직 시간대별 완료 파일이 없습니다.")
        return

    st.caption(
        f"완료 파일 {len(batches)}개 · 완료 파일을 삭제하면 해당 사진으로 다시 생성할 수 있습니다."
    )

    for idx, batch in enumerate(batches):
        created_at = batch.get("created_at", "")
        source_files = [x for x in batch.get("source_files", []) if x]
        source_text = ", ".join(source_files)
        sheet_names = [x for x in batch.get("sheet_names", []) if x]
        batch_id = str(batch.get("batch_id", "") or "")

        col_info, col_download, col_delete = st.columns([4.8, 1.15, 0.85])
        with col_info:
            st.markdown(f"**{created_at} 완료**")
            reused_count = int(batch.get("reused_records_count", 0) or 0)
            record_text = f"신규 기록 {batch.get('new_records_count', 0)}건"
            if reused_count:
                record_text += f" · 기존 기록 재사용 {reused_count}건"
            st.caption(
                f"업로드 파일 {batch.get('source_file_count', len(source_files))}개 · "
                f"{record_text} · 시트 {len(sheet_names)}개"
            )
            if source_text:
                st.caption(f"원본: {source_text}")

        with col_download:
            try:
                with open(batch["snapshot_path"], "rb") as f:
                    snapshot_bytes = f.read()
                st.download_button(
                    "다운로드",
                    data=snapshot_bytes,
                    file_name=batch.get("snapshot_filename", "체감온도측정_누적대장.xlsx"),
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                    key=f"heat_history_download_{batch_id or idx}",
                )
            except Exception as e:
                st.caption(f"파일 읽기 실패: {e}")

        with col_delete:
            if st.button(
                "삭제",
                use_container_width=True,
                key=f"heat_history_delete_{batch_id or idx}",
            ):
                ok, message = delete_heat_completed_batch(batch_id)
                if ok:
                    # 화면에 남아 있는 기존 업로드 캐시도 지워 재업로드가 즉시 가능하게 함.
                    st.session_state.pop("heat_saved", None)
                    st.session_state["heat_uploader_version"] = (
                        int(st.session_state.get("heat_uploader_version", 0)) + 1
                    )
                    st.rerun()
                else:
                    st.error(message)

        if idx < len(batches) - 1:
            st.markdown("---")


def save_heat_measurement(location: str, row: dict) -> Tuple[str, int, str, bool]:
    """측정 1건 저장. 같은 기록이면 새 행을 만들지 않고 기존 행을 반환."""
    template_ws = load_heat_template_sheet()
    if template_ws is None:
        raise ValueError(
            f"템플릿 파일을 찾을 수 없습니다: {HEAT_TEMPLATE_XLSX} "
            "(templates 폴더에 heat_index_template.xlsx를 추가해주세요)"
        )

    with _HEAT_HISTORY_LOCK:
        if os.path.exists(HEAT_LOG_FILE):
            wb = openpyxl.load_workbook(HEAT_LOG_FILE)
        else:
            wb = openpyxl.Workbook()
            wb.remove(wb.active)

        ws = get_or_create_heat_sheet(wb, template_ws, location, row["측정일자"])

        # 기존 누적 대장 자체도 검사해, 파일명이 달라진 같은 사진/같은 측정값의 중복 저장을 막음.
        duplicate_row = find_duplicate_heat_row(ws, row)
        if duplicate_row is not None:
            existing_note = str(ws.cell(row=duplicate_row, column=8).value or "").strip()
            duplicate_note = "중복 기록 건너뜀(기존 기록 유지)"
            display_note = " / ".join([x for x in [existing_note, duplicate_note] if x])
            return ws.title, duplicate_row, display_note, True

        target_row = find_empty_heat_slot(ws, template_ws)
        if target_row is None:
            raise ValueError(f"'{ws.title}' 기록은 이미 9건이 모두 채워져 있습니다.")

        notes = []
        row["측정자"] = format_measurer(row.get("측정자", ""))
        if row.get("_measurer_unrecognized"):
            notes.append("⚠️측정자 확인필요(등록된 11명 외 이름은 저장하지 않음)")

        # 판독 공백 자동 보완: 같은 장소·날짜 평균을 우선 사용하고,
        # 자료가 없으면 같은 날짜의 다른 장소 평균을 보조로 사용한다.
        same_sheet_averages = get_heat_field_averages_from_sheet(ws)
        same_date_averages = get_heat_field_averages_from_workbook_date(
            wb, row["측정일자"], exclude_sheet=ws.title
        )
        _, _, unavailable_fields = fill_missing_heat_values(
            row, same_sheet_averages, same_date_averages
        )
        notes.extend(build_heat_auto_notes(row, unavailable_fields))

        new_min = heat_time_to_minutes(row["측정시간"])
        if target_row > 6 and new_min is not None:
            for prev_row in range(target_row - 1, 5, -1):
                prev_min = heat_time_to_minutes(ws.cell(row=prev_row, column=3).value)
                if prev_min is not None:
                    if new_min - prev_min > HEAT_LOG_GAP_MINUTES:
                        notes.append(f"간격초과({new_min - prev_min}분 경과, 측정 누락 가능성)")
                    break

        if is_implausible_value(row["기온"], row["습도"]):
            notes.append("⚠️확인필요(비현실적 수치, OCR 오독 가능성 - 직접 확인 후 수정 필요)")

        gap_note = " / ".join(notes)

        ws.cell(row=target_row, column=3, value=format_heat_time_display(row["측정시간"]))
        ws.cell(row=target_row, column=4, value=_to_number(row["기온"]))
        ws.cell(row=target_row, column=5, value=_to_number(row["습도"]))
        ws.cell(row=target_row, column=6, value=_to_number(row["체감온도"]))
        ws.cell(row=target_row, column=7, value=format_measurer(row["측정자"]))
        ws.cell(row=target_row, column=8, value=gap_note)

        wb.save(HEAT_LOG_FILE)
        return ws.title, target_row, gap_note, False


def overwrite_heat_row(sheet_name: str, target_row: int, fields: dict) -> Tuple[dict, str]:
    """이미 저장된 칸을 덮어씀. 수정 화면에서 값을 비워도 주변 평균으로 다시 보완."""
    if not os.path.exists(HEAT_LOG_FILE):
        raise ValueError("저장된 기록 파일이 없습니다.")
    wb = openpyxl.load_workbook(HEAT_LOG_FILE)
    if sheet_name not in wb.sheetnames:
        raise ValueError(f"'{sheet_name}' 시트를 찾을 수 없습니다.")
    ws = wb[sheet_name]

    # 측정자는 지정된 11명만 저장하며 직책과 관계없이 이름 대원 형식으로 통일한다.
    raw_person = str(fields.get("측정자", "") or "").strip()
    normalized_person = format_measurer(raw_person)
    if raw_person and not normalized_person:
        raise ValueError("측정자는 지정된 11명 중 한 명만 입력할 수 있습니다.")
    fields["측정자"] = normalized_person

    # 사용자가 수정 입력에서 빈 값을 남긴 경우에도 같은 장소·날짜 평균으로 보완한다.
    fields.pop("_average_filled", None)
    fields.pop("_heat_index_calculated", None)
    same_sheet_averages = get_heat_field_averages_from_sheet(ws, exclude_row=target_row)

    date_match = re.search(r"_(\d{4})(?:_\d+)?$", sheet_name)
    date_str = date_match.group(1) if date_match else ""
    same_date_averages = get_heat_field_averages_from_workbook_date(
        wb, date_str, exclude_sheet=sheet_name
    )
    _, _, unavailable_fields = fill_missing_heat_values(
        fields, same_sheet_averages, same_date_averages
    )

    existing_notes = clean_heat_auto_notes(ws.cell(row=target_row, column=8).value)
    existing_notes.extend(build_heat_auto_notes(fields, unavailable_fields))
    note_text = " / ".join(existing_notes)

    ws.cell(row=target_row, column=3, value=format_heat_time_display(fields["측정시간"]))
    ws.cell(row=target_row, column=4, value=_to_number(fields["기온"]))
    ws.cell(row=target_row, column=5, value=_to_number(fields["습도"]))
    ws.cell(row=target_row, column=6, value=_to_number(fields["체감온도"]))
    ws.cell(row=target_row, column=7, value=format_measurer(fields["측정자"]))
    ws.cell(row=target_row, column=8, value=note_text)
    wb.save(HEAT_LOG_FILE)
    return fields, note_text


def render_heat_index_log():
    st.markdown("---")
    st.markdown(
        """
        <h2 style="color:#e8590c; font-weight:800; font-size:48px; margin-top:0.35rem; margin-bottom:0.35rem;">
            체감온도 측정 기록
        </h2>
        """,
        unsafe_allow_html=True
    )

    if openpyxl is None:
        st.error("openpyxl 패키지가 설치되어 있지 않습니다. requirements.txt에 openpyxl을 추가해주세요.")
        return

    if not os.path.exists(HEAT_TEMPLATE_XLSX):
        st.error(
            "템플릿 파일이 없습니다. 저장소의 templates/heat_index_template.xlsx 위치에 "
            "체감온도측정 대장 양식을 추가해주세요."
        )
        return

    if "GPT_API_KEY" not in st.secrets:
        st.error("Secrets에 GPT_API_KEY 설정이 필요합니다. (기존 번역 기능과 같은 키를 사용합니다)")
        return

    api_key = st.secrets["GPT_API_KEY"]

    if "heat_saved" not in st.session_state:
        st.session_state["heat_saved"] = {}
    if "heat_uploader_version" not in st.session_state:
        st.session_state["heat_uploader_version"] = 0

    heat_files = st.file_uploader(
        "측정 사진 업로드",
        accept_multiple_files=True,
        type=["jpg", "png", "jpeg", "webp", "heic", "heif"],
        key=f"heat_index_uploader_{st.session_state['heat_uploader_version']}"
    )

    new_file_infos = []
    new_entries_for_snapshot = []
    render_items = []

    if heat_files:
        for f in heat_files:
            file_hash, file_bytes = uploaded_file_hash(f)
            session_key = file_hash

            if session_key in st.session_state["heat_saved"]:
                render_items.append(st.session_state["heat_saved"][session_key])
                continue

            should_process, existing_meta = reserve_heat_upload(file_hash, f.name, len(file_bytes))

            if not should_process:
                item = {
                    "file_hash": file_hash,
                    "file_name": f.name,
                    "duplicate_file": True,
                    "meta": existing_meta,
                    "entries": [],
                }
                st.session_state["heat_saved"][session_key] = item
                render_items.append(item)
                continue

            suffix = os.path.splitext(f.name)[1].lower() or ".jpg"
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(file_bytes)
                original_path = tmp.name

            entries = []
            process_error = ""
            with st.spinner(f"{f.name} 분석 및 자동저장 중..."):
                try:
                    records = extract_heat_records_with_gpt(api_key, original_path)
                    records = prepare_heat_records_for_average(records)

                    for rec in records:
                        try:
                            sheet_name, target_row, gap_note, was_duplicate = save_heat_measurement(
                                rec["측정위치"], rec
                            )
                            rec["측정시간"] = format_heat_time_display(rec["측정시간"])
                            rec["측정자"] = format_measurer(rec["측정자"])
                            entries.append({
                                "sheet": sheet_name,
                                "row": target_row,
                                "gap_note": gap_note,
                                "values": rec,
                                "duplicate": was_duplicate,
                                "error": None,
                            })
                        except Exception as e:
                            entries.append({
                                "sheet": None,
                                "row": None,
                                "gap_note": "",
                                "values": rec,
                                "duplicate": False,
                                "error": str(e),
                            })

                    status = "completed" if records else "no_records"
                    finalize_heat_upload(file_hash, status, entries)

                except Exception as e:
                    process_error = str(e)
                    finalize_heat_upload(file_hash, "error", entries, error=process_error)

            if os.path.exists(original_path):
                try:
                    os.remove(original_path)
                except Exception:
                    pass

            item = {
                "file_hash": file_hash,
                "file_name": f.name,
                "duplicate_file": False,
                "meta": {
                    "status": "error" if process_error else ("completed" if entries else "no_records"),
                    "error": process_error,
                },
                "entries": entries,
            }
            st.session_state["heat_saved"][session_key] = item
            render_items.append(item)

            if not process_error:
                new_file_infos.append({
                    "file_hash": file_hash,
                    "original_name": f.name,
                })
                new_entries_for_snapshot.extend(entries)

        # 한 번에 선택한 새 사진들을 하나의 시간대 완료 파일로 묶음.
        created_batch = create_heat_batch_snapshot(new_file_infos, new_entries_for_snapshot)
        if created_batch:
            reused_count = int(created_batch.get("reused_records_count", 0) or 0)
            result_text = f"신규 기록 {created_batch.get('new_records_count', 0)}건"
            if reused_count:
                result_text += f" · 기존 기록 재사용 {reused_count}건"
            st.success(
                f"시간대별 완료 파일 생성: {created_batch['created_at']} · {result_text}"
            )

        for item in render_items:
            file_name = item.get("file_name", "업로드 파일")

            if item.get("duplicate_file"):
                meta = item.get("meta", {})
                processed_at = meta.get("processed_at") or meta.get("started_at") or "이전 처리 시각"
                records_count = meta.get("records_count", 0)
                with st.expander(f"♻️ {file_name} — 이미 처리된 파일", expanded=False):
                    st.info(
                        f"동일한 파일 내용이 {processed_at}에 이미 처리되어 새 행을 만들지 않았습니다. "
                        f"기존 인식 기록: {records_count}건"
                    )
                    if os.path.exists(HEAT_LOG_FILE):
                        with open(HEAT_LOG_FILE, "rb") as f:
                            duplicate_master_bytes = f.read()
                        st.download_button(
                            "현재 누적 대장 엑셀 다운로드",
                            data=duplicate_master_bytes,
                            file_name=f"체감온도측정_전체누적대장_{datetime.now().strftime('%Y%m%d')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            key=f"heat_duplicate_master_download_{file_hash}",
                        )
                continue

            entries = item.get("entries", [])
            meta = item.get("meta", {})
            process_error = meta.get("error", "")

            with st.expander(f"📷 {file_name} — {len(entries)}건 인식", expanded=True):
                if process_error:
                    st.error(f"분석 실패: {process_error}")
                    continue

                if not entries:
                    st.warning("이 사진에서 측정 기록을 찾지 못했습니다. 같은 사진의 자동 반복 분석은 막았습니다.")

                for i, entry in enumerate(entries):
                    if entry["error"]:
                        st.error(f"저장 실패: {entry['error']}")
                        continue

                    if entry.get("duplicate"):
                        st.info(
                            f"♻️ 기존 기록과 동일하여 새 행을 만들지 않음 → "
                            f"'{entry['sheet']}' 시트 {entry['row']-5}번째 줄 유지"
                        )
                    else:
                        st.success(f"✅ 자동저장됨 → '{entry['sheet']}' 시트 {entry['row']-5}번째 줄")

                    if "⚠️확인필요" in entry["gap_note"]:
                        st.error(entry["gap_note"])
                    elif "간격초과" in entry["gap_note"]:
                        st.warning(entry["gap_note"])
                    elif entry["gap_note"] and not entry.get("duplicate"):
                        st.info(entry["gap_note"])

                    wb_preview = openpyxl.load_workbook(HEAT_LOG_FILE)
                    if entry["sheet"] in wb_preview.sheetnames:
                        st.caption(
                            f"오늘 이 장소 측정현황(4회 기준): "
                            f"{get_slot_coverage_text(wb_preview[entry['sheet']])}"
                        )

                    # 중복으로 건너뛴 기존 행은 현재 업로드에서 수정하지 않음.
                    if entry.get("duplicate"):
                        continue

                    v = entry["values"]
                    ec1, ec2 = st.columns(2)
                    with ec1:
                        st.caption(f"측정위치: {v['측정위치']}  |  측정일자: {v['측정일자']}")
                        new_time = st.text_input(
                            "측정시간", v["측정시간"],
                            key=f"heat_time_{item['file_hash']}_{i}"
                        )
                        new_temp = st.text_input(
                            "기온", v["기온"],
                            key=f"heat_temp_{item['file_hash']}_{i}"
                        )
                    with ec2:
                        new_hum = st.text_input(
                            "습도", v["습도"],
                            key=f"heat_hum_{item['file_hash']}_{i}"
                        )
                        new_feels = st.text_input(
                            "체감온도", v["체감온도"],
                            key=f"heat_feels_{item['file_hash']}_{i}"
                        )
                        new_person = st.text_input(
                            "측정자", v["측정자"],
                            key=f"heat_person_{item['file_hash']}_{i}"
                        )

                    if st.button("수정 반영", key=f"heat_fix_{item['file_hash']}_{i}"):
                        fixed = {
                            "측정시간": new_time.strip(),
                            "기온": new_temp.strip(),
                            "습도": new_hum.strip(),
                            "체감온도": new_feels.strip(),
                            "측정자": new_person.strip(),
                        }
                        try:
                            fixed, updated_note = overwrite_heat_row(entry["sheet"], entry["row"], fixed)
                            fixed["측정시간"] = format_heat_time_display(fixed["측정시간"])
                            fixed["측정자"] = format_measurer(fixed["측정자"])
                            entry["values"].update(fixed)
                            entry["gap_note"] = updated_note
                            st.success("수정 반영되었습니다. 빈 측정값은 평균치로 자동 보완됩니다.")
                        except Exception as e:
                            st.error(f"수정 실패: {e}")

    st.markdown("#### 누적 기록")
    if os.path.exists(HEAT_LOG_FILE):
        wb = openpyxl.load_workbook(HEAT_LOG_FILE)
        st.caption(f"저장된 일자별 기록 시트: {', '.join(wb.sheetnames)}")

        today_str = datetime.now().strftime("%Y%m%d")
        with open(HEAT_LOG_FILE, "rb") as f:
            master_bytes = f.read()
        st.download_button(
            "전체 누적 대장 엑셀 다운로드",
            data=master_bytes,
            file_name=f"체감온도측정_전체누적대장_{today_str}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key="heat_log_download_btn"
        )
    else:
        st.info("아직 저장된 측정 기록이 없습니다.")

    render_heat_completed_file_list()


def main():
    st.set_page_config(page_title="TBM PPT Maker", layout="wide")
    hide_streamlit_ui()

    render_app_title()

    render_daily_safety_meeting()

    st.markdown("---")
    st.markdown(
        """
        <h2 style="color:#0057d9; font-weight:800; font-size:48px; margin-top:0.35rem; margin-bottom:0.35rem;">
            TBM 번역 PPT
        </h2>
        """,
        unsafe_allow_html=True
    )

    render_tbm_input_area()

    render_shared_notice_board()

    render_heat_index_log()


if __name__ == "__main__":
    main()